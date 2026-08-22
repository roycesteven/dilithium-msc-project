#!/usr/bin/env python3
"""Render the 6-8 minute video deck into the University of Manchester template.

This is the PowerPoint twin of scripts/gen_slides.py.  Both read the SAME two
sources, so the .pptx, the HTML deck and the report can never quote different
evidence:

  * report/latex/generated/*.tex   -- every number, via gen_slides.read_macros()
  * report/slides/video_deck.template.html
        -- the slide order, the slide titles, the planned times and the speaker
           notes, taken from the rendered template rather than retyped

Brand chrome comes from `report/slides/Master_169 presentation(2).pptx`: the
deck is BUILT ON that file, so the slide master, theme, Arial type scheme, the
purple #7800A2 title colour, the #660066 dotted rule and the Manchester logo
(a master shape, inherited by every slide) are the template's own, not a
reconstruction of them.

The drawn scenes are rebuilt as NATIVE PowerPoint shapes -- rectangles,
connectors and text boxes -- so every element stays selectable and editable in
place.  Only the two report figures (fig_timing, fig_onchain) are pictures,
because they are the report's own committed figures and must stay identical to
what the report prints.

This script only READS evidence and report artefacts.  It never builds, runs a
benchmark, or estimates a value; a missing macro is a hard error rather than a
blank on a slide.

Requires python-pptx (pip install python-pptx); nothing else beyond what
gen_slides.py already needs (poppler's pdftoppm, to rasterise fig_timing).

    python3 scripts/gen_slides_pptx.py
        -> report/slides/video_deck_uom.pptx

Check the layout by RENDERING, never by reading the code: PowerPoint is
reachable from WSL, so the deck can be exported to PDF and looked at --

    cp report/slides/video_deck_uom.pptx /mnt/c/Users/<user>/AppData/Local/Temp/deck.pptx
    powershell.exe -NoProfile -Command "\
      $pp = New-Object -ComObject PowerPoint.Application; \
      $d  = $pp.Presentations.Open('C:\\Users\\<user>\\AppData\\Local\\Temp\\deck.pptx', \
                                   $true, $false, $false); \
      $d.SaveAs('C:\\Users\\<user>\\AppData\\Local\\Temp\\deck.pdf', 32); \
      $d.Close(); $pp.Quit()"
"""

from __future__ import annotations

import argparse
import html
import pathlib
import re
import shutil
import subprocess
import sys
import tempfile

try:
    from pptx import Presentation
except ModuleNotFoundError:                                   # pragma: no cover
    sys.exit("python-pptx is not installed -- pip install python-pptx")
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = pathlib.Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "scripts"))
import gen_slides  # noqa: E402  (same macro reader as the HTML deck)

SLIDES = ROOT / "report" / "slides"
TEMPLATE_HTML = SLIDES / "video_deck.template.html"
UOM_PPTX = SLIDES / "Master_169 presentation(2).pptx"
OUTPUT = SLIDES / "video_deck_uom.pptx"
FIG_DIR = ROOT / "report" / "latex" / "figures"
FIG_DPI = 200

# ---------------------------------------------------------------- palette ---
# Template chrome, read out of the UoM master (never invented):
UOM_PURPLE = RGBColor(0x78, 0x00, 0xA2)   # slideMaster1 titleStyle
UOM_RULE = RGBColor(0x66, 0x00, 0x66)     # layout1 straight connector, dotted
UOM_GREY = RGBColor(0x59, 0x59, 0x59)     # template title/sub-title text

# The deck's semantic accents, carried over unchanged: they mean "basic
# signature" / "adaptor layer" / "reused" / "warning" throughout the report and
# in the two embedded report figures, so recolouring them would break the tie
# between a slide and the figure printed beside it.
INK = RGBColor(0x14, 0x21, 0x3A)
MUTED = RGBColor(0x5A, 0x6B, 0x86)
LINE = RGBColor(0xD8, 0xE0, 0xEC)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xEE, 0xF3, 0xFA)
BASE = RGBColor(0x1F, 0x6F, 0xEB)
BASE_SOFT = RGBColor(0xE8, 0xF1, 0xFF)
ADAPTOR = RGBColor(0xE8, 0x59, 0x0C)
ADAPTOR_SOFT = RGBColor(0xFD, 0xEC, 0xE1)
INK_AD = RGBColor(0xC2, 0x41, 0x0C)
REUSED = RGBColor(0x2F, 0x9E, 0x44)
REUSED_SOFT = RGBColor(0xE6, 0xF6, 0xEA)
WARN = RGBColor(0xB0, 0x2A, 0x37)
WARN_SOFT = RGBColor(0xFD, 0xEC, 0xEE)

SANS = "Arial"          # the template theme's major AND minor latin typeface
MONO = "Consolas"

# ---------------------------------------------------------------- geometry --
SW, SH = 12192000, 6858000
L = 523875                       # master logo's own left edge
R = 11668125                     # mirror margin
W = R - L
LOGO_BOTTOM = 1220788            # master picture: 509588 + 711200

STEP_Y = 640000                  # "where are we" stepper, right of the logo
EYE_Y = 1295000                  # eyebrow, clear of the logo
TTL_TOP, TTL_BOT = 1495000, 2105000
RULE_Y = 2165000
BODY_Y = 2265000
BODY_B = 6430000
BODY_H = BODY_B - BODY_Y

# Headline sizing.  The title band is fixed (the logo above it and the purple
# rule below it are both template chrome), so the SIZE has to follow the string
# or a two-line headline climbs into the eyebrow -- which is exactly what the
# first render did on three slides.  Arial bold averages ~0.56 em, and the band
# holds two lines at 22 pt.
TITLE_STEPS = ((62, 25.0), (140, 22.0), (10**6, 19.0))

GAP = 170000                     # the deck's 16px column gap, to scale
PARTS = [("why", "Why"), ("method", "Method"), ("results", "Results"),
         ("takeaway", "Takeaway")]


# ============================================================== primitives ==
def emu(v: float) -> int:
    return int(round(v))


def no_shadow(shape) -> None:
    try:
        shape.shadow.inherit = False
    except (AttributeError, NotImplementedError):
        pass


def set_spacing(run, hundredths_pt: int) -> None:
    """Letter-spacing: the deck's eyebrow/label tracking, which python-pptx
    does not expose as a property."""
    run.font._rPr.set("spc", str(hundredths_pt))


MARKUP = re.compile(r"\[\[(\w+)\|([^\]]*)\]\]")
STYLES = {
    "b": dict(bold=True, color=INK),
    "hi": dict(bold=True, color=INK_AD),
    "hb": dict(bold=True, color=BASE),
    "ok": dict(bold=True, color=REUSED),
    "warn": dict(bold=True, color=WARN),
    "m": dict(mono=True),
    "mad": dict(mono=True, color=INK_AD),
    "mu": dict(color=MUTED),
    "ink": dict(color=INK),
    "purple": dict(bold=True, color=UOM_PURPLE),
}


def split_markup(text: str):
    """[[style|text]] -> (chunk, style-name-or-None) pairs."""
    out, pos = [], 0
    for m in MARKUP.finditer(text):
        if m.start() > pos:
            out.append((text[pos:m.start()], None))
        out.append((m.group(2), m.group(1)))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], None))
    return out


def write(tf, text, *, size, color=INK, bold=False, italic=False, mono=False,
          align=PP_ALIGN.LEFT, line=1.25, space_before=0, space_after=0,
          caps=False, spacing=None, first=False):
    """Append one paragraph, honouring the [[style|...]] inline markup."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    for chunk, style in split_markup(text):
        if not chunk:
            continue
        s = STYLES.get(style, {}) if style else {}
        run = p.add_run()
        run.text = chunk.upper() if caps else chunk
        f = run.font
        f.name = MONO if s.get("mono", mono) else SANS
        f.size = Pt(s.get("size", size))
        f.bold = s.get("bold", bold)
        f.italic = italic
        f.color.rgb = s.get("color", color)
        if spacing:
            set_spacing(run, spacing)
    return p


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    no_shadow(box)
    return box, tf


def para_box(slide, x, y, w, h, text, **kw):
    """One-paragraph text box (the common case)."""
    anchor = kw.pop("anchor", MSO_ANCHOR.TOP)
    box, tf = textbox(slide, x, y, w, h, anchor=anchor, wrap=kw.pop("wrap", True))
    write(tf, text, first=True, **kw)
    return box, tf


def rect(slide, x, y, w, h, *, fill=None, stroke=None, stroke_w=12700,
         radius=None, dash=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, emu(x), emu(y), emu(w), emu(h))
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = min(0.5, radius / max(1.0, min(w, h)))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if stroke is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = stroke
        sp.line.width = emu(stroke_w)
        if dash:
            sp.line.dash_style = dash
    no_shadow(sp)
    sp.text_frame.word_wrap = True
    return sp


def card(slide, x, y, w, h, *, fill=CARD, stroke=LINE, radius=110000):
    sp = rect(slide, x, y, w, h, fill=fill, stroke=stroke, stroke_w=12700,
              radius=radius, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = Emu(160000)
    tf.margin_top = tf.margin_bottom = Emu(120000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return sp, tf


def line(slide, x1, y1, x2, y2, *, color=LINE, width=12700, dash=None,
         arrow=False, arrow_start=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      emu(x1), emu(y1), emu(x2), emu(y2))
    conn.line.color.rgb = color
    conn.line.width = emu(width)
    if dash:
        conn.line.dash_style = dash
    if arrow or arrow_start:
        ln = conn.line._get_or_add_ln()
        if arrow_start:
            head = ln.makeelement(qn("a:headEnd"), {})
            head.set("type", "triangle")
            head.set("w", "med")
            head.set("len", "med")
            ln.append(head)
        if arrow:
            tail = ln.makeelement(qn("a:tailEnd"), {})
            tail.set("type", "triangle")
            tail.set("w", "med")
            tail.set("len", "med")
            ln.append(tail)
    no_shadow(conn)
    return conn


def chip(slide, x, y, text, *, color=MUTED, border=LINE, fill=None, size=9.0,
         mono=True, height=250000, pad=110000):
    """Pill label; width follows the string so a row of chips can be packed."""
    advance = size * (0.55 if mono else 0.52) * 12700
    w = len(text) * advance + 2 * pad
    sp = rect(slide, x, y, w, height, fill=fill, stroke=border, stroke_w=9525,
              radius=height / 2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = Emu(emu(pad))
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    write(tf, text, size=size, color=color, mono=mono, first=True,
          align=PP_ALIGN.CENTER, line=1.0)
    return sp, w


def chipline(slide, x, y, items, *, gap=90000, height=250000, size=9.0):
    cx = x
    for text, colour in items:
        _, w = chip(slide, cx, y, text, color=colour, border=colour,
                    size=size, height=height)
        cx += w + gap
    return cx


def caveat(slide, x, y, w, h, text, *, size=10.0, color=MUTED, line_h=1.3):
    """The deck's caveat block: a warn-coloured left rule, then muted text."""
    rect(slide, x, y, 38000, h, fill=WARN, stroke=None)
    para_box(slide, x + 150000, y, w - 150000, h, text,
             size=size, color=color, line=line_h)


# ----------------------------------------------------------------- scenes ---
class Scene:
    """Maps an SVG viewBox onto the slide so the deck's drawn scenes can be
    transcribed coordinate-for-coordinate into native shapes."""

    def __init__(self, slide, x, y, w, vb):
        self.slide = slide
        self.x, self.y = x, y
        self.vbx, self.vby, self.vw, self.vh = vb
        self.u = w / self.vw
        self.w = w
        self.h = self.vh * self.u

    # coordinate helpers -------------------------------------------------
    def X(self, v):
        return self.x + (v - self.vbx) * self.u

    def Y(self, v):
        return self.y + (v - self.vby) * self.u

    def U(self, v):
        return v * self.u

    def pt(self, px):
        return px * self.u / 12700.0

    # drawing ------------------------------------------------------------
    def rect(self, x, y, w, h, **kw):
        radius = kw.pop("r", None)
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        return rect(self.slide, self.X(x), self.Y(y), self.U(w), self.U(h),
                    radius=self.U(radius) if radius else None, shape=shape,
                    stroke_w=kw.pop("sw", 1.4) * self.u, **kw)

    def text(self, x, y, px, text, *, weight=False, fill=INK, anchor="start",
             mono=False, spacing=None, width=600):
        size = self.pt(px)
        top = self.Y(y) - size * 12700 * 0.86
        h = size * 12700 * 1.45
        if anchor == "middle":
            left, align = self.X(x) - self.U(width) / 2, PP_ALIGN.CENTER
        elif anchor == "end":
            left, align = self.X(x) - self.U(width), PP_ALIGN.RIGHT
        else:
            left, align = self.X(x), PP_ALIGN.LEFT
        box, tf = textbox(self.slide, left, top, self.U(width), h, wrap=False)
        write(tf, text, size=size, color=fill, bold=weight, mono=mono,
              align=align, line=1.0, spacing=spacing, first=True)
        return box

    def line(self, x1, y1, x2, y2, **kw):
        kw.setdefault("width", 2.0 * self.u)
        return line(self.slide, self.X(x1), self.Y(y1), self.X(x2), self.Y(y2),
                    **kw)


# ------------------------------------------------------------------ input ---
SECTION = re.compile(r'<section class="slide[^"]*"([^>]*)>')
ATTR = re.compile(r'data-(title|time|notes|part)="([^"]*)"')


def deck_meta(values: dict[str, str]) -> list[dict[str, str]]:
    """Slide titles, planned times, parts and speaker notes -- rendered from
    the same template + macros the HTML deck is rendered from."""
    # Only the text attributes are wanted here, so the embedded images are
    # stubbed out -- but every placeholder the template uses must still resolve,
    # or render() (rightly) refuses.
    stubs = dict(values, FIG_TIMING="", FIG_ONCHAIN="", UOM_LOGO="")
    rendered = gen_slides.render(TEMPLATE_HTML, stubs)
    out = []
    for m in SECTION.finditer(rendered):
        attrs = {k: html.unescape(v) for k, v in ATTR.findall(m.group(1))}
        out.append(attrs)
    return out


def figure_png(stem: str, tmpdir: pathlib.Path) -> pathlib.Path:
    png = FIG_DIR / f"{stem}.png"
    if png.is_file():
        return png
    pdf = FIG_DIR / f"{stem}.pdf"
    if not pdf.is_file():
        sys.exit(f"missing figure {pdf}")
    if not shutil.which("pdftoppm"):
        sys.exit("pdftoppm not found (poppler-utils) -- needed for the figures")
    out = tmpdir / stem
    subprocess.run(["pdftoppm", "-png", "-r", str(FIG_DPI), "-singlefile",
                    str(pdf), str(out)], check=True, capture_output=True)
    return out.with_suffix(".png")


def picture_fit(slide, path, x, y, w, h):
    """Contain-fit a picture inside a box, centred (the deck's .figbox)."""
    pic = slide.shapes.add_picture(str(path), emu(x), emu(y))
    scale = min(w / pic.width, h / pic.height)
    pic.width = emu(pic.width * scale)
    pic.height = emu(pic.height * scale)
    pic.left = emu(x + (w - pic.width) / 2)
    pic.top = emu(y + (h - pic.height) / 2)
    return pic


# =============================================================== chrome =====
def new_slide(prs, meta, index, total, *, chrome=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # Blank: logo from master
    if chrome:
        part = meta.get("part", "")
        if part:
            cx = R
            for key, label in reversed(PARTS):
                on = key == part
                _, w = chip(slide, 0, STEP_Y, label,
                            color=INK_AD if on else MUTED,
                            border=ADAPTOR if on else LINE,
                            fill=ADAPTOR_SOFT if on else None,
                            size=8.5, mono=True, height=230000, pad=95000)
                cx -= w
                slide.shapes[-1].left = emu(cx)
                cx -= 40000
        para_box(slide, R - 900000, BODY_B + 120000, 900000, 220000,
                 f"{index} / {total}", size=8.5, color=MUTED, mono=True,
                 align=PP_ALIGN.RIGHT)
    notes = meta.get("notes", "")
    planned = meta.get("time", "")
    tf = slide.notes_slide.notes_text_frame
    tf.text = f"[planned {planned}]  {meta.get('title','')}"
    p = tf.add_paragraph()
    p.text = notes
    return slide


def head(slide, eyebrow, title, *, title_size=None):
    plain = "".join(chunk for chunk, _ in split_markup(title))
    if title_size is None:
        title_size = next(pt for limit, pt in TITLE_STEPS if len(plain) <= limit)
    para_box(slide, L, EYE_Y, W, 210000, eyebrow, size=9.5, color=MUTED,
             mono=True, caps=True, spacing=170)
    para_box(slide, L, TTL_TOP, W, TTL_BOT - TTL_TOP, title,
             size=title_size, color=UOM_PURPLE, bold=True, line=1.08,
             anchor=MSO_ANCHOR.BOTTOM)
    line(slide, L, RULE_Y, R, RULE_Y, color=UOM_RULE, width=25400,
         dash=MSO_LINE_DASH_STYLE.ROUND_DOT)


# ================================================================ slides ====
def slide_title(prs, V, meta, idx, total):
    """Slide 1 -- the template's own Title Slide layout, filled."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)

    para_box(slide, L, 1560000, W, 240000,
             "MSc Cybersecurity · University of Manchester · COMP66060",
             size=10.5, color=MUTED, mono=True, caps=True, spacing=170)
    para_box(slide, L, 1850000, W, 1500000,
             "Implementing Post-Quantum Secure\n"
             "Exotic Signature Schemes in Blockchains",
             size=30.0, color=UOM_GREY, bold=True, line=1.16)
    line(slide, L, 3789040, L + 9505056, 3789040, color=UOM_RULE, width=25400,
         dash=MSO_LINE_DASH_STYLE.ROUND_DOT)
    para_box(slide, L, 4180000, 9200000, 900000,
             "LAS — a lattice adaptor signature built on CRYSTALS-Dilithium "
             "primitives: implemented twice, benchmarked against two baselines, "
             "and demonstrated end-to-end in a post-quantum atomic swap.",
             size=15.0, color=UOM_GREY, line=1.3)
    para_box(slide, L, 5300000, W, 300000,
             f"Royce Steven     ·     Supervisor: Dr Zhipeng Wang     ·     "
             f"[[m|evidence run {V['benchRunId']}]]",
             size=11.5, color=MUTED)

    tf = slide.notes_slide.notes_text_frame
    tf.text = f"[planned {meta.get('time','')}]  {meta.get('title','')}"
    tf.add_paragraph().text = meta.get("notes", "")
    return slide


def slide_why(prs, V, meta, idx, total):
    """Slide 2 -- the application, drawn."""
    s = new_slide(prs, meta, idx, total)
    head(s, "The application first",
         "Two strangers, two chains, [[hb|no trusted middleman]]")

    sc = Scene(s, L, BODY_Y, W, (0, 44, 1160, 252))
    # the two parties
    for x0, name, sub, coin, coin_col, chain, chain_no, foot in (
        (1, "Alice", "has bitcoin, wants ether", "1 BTC", RGBColor(0xF7, 0x93, 0x1A),
         "Bitcoin", "chain 1", "spendable only by her own key"),
        (893, "Bob", "has ether, wants bitcoin", "10 ETH", RGBColor(0x62, 0x7E, 0xEA),
         "Ethereum", "chain 2", "spendable only by his own key"),
    ):
        sc.rect(x0, 44, 266, 252, r=14, fill=CARD, stroke=LINE, sw=1.5)
        sc.text(x0 + 25, 86, 23, name, weight=True, width=200)
        sc.text(x0 + 25, 110, 15, sub, fill=MUTED, width=250)
        sc.rect(x0 + 29, 156, 72, 72, r=36, fill=None, stroke=coin_col, sw=2.2)
        sc.text(x0 + 65, 198, 15, coin, weight=True, anchor="middle", width=200)
        sc.rect(x0 + 117, 166, 126, 52, r=11, fill=WASH, stroke=LINE, sw=1.2)
        sc.text(x0 + 133, 188, 14, chain, weight=True, width=150)
        sc.text(x0 + 133, 207, 12.5, chain_no, fill=MUTED, width=150)
        sc.text(x0 + 25, 266, 13, foot, fill=MUTED, width=260)

    # the two payments
    orange = RGBColor(0xF7, 0x93, 0x1A)
    blue = RGBColor(0x62, 0x7E, 0xEA)
    sc.text(580, 104, 17, "Alice’s bitcoin pays Bob", weight=True,
            anchor="middle", width=500)
    sc.line(296, 126, 866, 126, color=orange, width=7 * sc.u, arrow=True)
    sc.line(864, 250, 294, 250, color=blue, width=7 * sc.u, arrow=True)
    sc.text(580, 288, 17, "Bob’s ether pays Alice", weight=True,
            anchor="middle", width=500)

    # the shared secret
    sc.line(580, 132, 580, 152, color=ADAPTOR, width=1.8 * sc.u,
            dash=MSO_LINE_DASH_STYLE.DASH)
    sc.line(580, 224, 580, 244, color=ADAPTOR, width=1.8 * sc.u,
            dash=MSO_LINE_DASH_STYLE.DASH)
    sc.rect(466, 152, 228, 72, r=14, fill=ADAPTOR_SOFT, stroke=ADAPTOR, sw=1.8)
    sc.text(580, 186, 15.5, "one shared secret", weight=True, fill=INK_AD,
            anchor="middle", width=220)
    sc.text(580, 206, 13.5, "gates both payments", fill=INK_AD,
            anchor="middle", width=220)

    y = BODY_Y + sc.h + 130000
    para_box(s, L, y, W, 320000,
             "Only the holder of the secret can complete the first payment — "
             "and publishing it reveals the secret.",
             size=13.5, bold=True, align=PP_ALIGN.CENTER)
    para_box(s, L, y + 300000, W, 300000,
             "That reveal releases the other leg — no trusted intermediary, and "
             "either both transfers happen or neither does.",
             size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    caveat(s, L, 5420000, W, 620000,
           "[[warn|The adaptor signatures in use on blockchains today are "
           "elliptic-curve]] — which Shor’s algorithm breaks. NIST has "
           "standardised [[warn|basic]] post-quantum signatures; the exotic "
           "ones [[warn|remain unevenly implemented]] — multi-signatures are "
           "actively being built, and that is the gap this project addresses "
           "for the adaptor case.", size=11.0)
    para_box(s, L, 6110000, W, 330000,
             "So: [[hi|can we build one]] — [[hi|what does it cost]] — and "
             "[[hi|will a real chain take it?]]", size=15.0, bold=True)
    return s


def slide_mechanism(prs, V, meta, idx, total):
    """Slide 3 -- the four functions, drawn."""
    s = new_slide(prs, meta, idx, total)
    head(s, "The mechanism · four functions added to a plain signature",
         "Publishing a signature [[hi|leaks a secret]] — that is what makes a "
         "swap atomic")

    sc = Scene(s, L, BODY_Y, W, (0, 46, 1160, 354))
    sc.rect(1, 56, 318, 152, r=13, fill=CARD, stroke=ADAPTOR, sw=1.6,
            dash=MSO_LINE_DASH_STYLE.DASH)
    sc.text(21, 86, 12, "1 · PRESIGN", weight=True, fill=ADAPTOR, spacing=110,
            width=250)
    sc.text(21, 118, 17, "σ̂ = (c, ẑ)", mono=True, width=250)
    sc.text(21, 144, 13.5, "accepted by PreVerify against Y", fill=MUTED,
            width=300)
    sc.rect(21, 158, 278, 36, r=9, fill=WARN_SOFT, stroke=WARN, sw=1.2)
    sc.text(35, 181, 13, "not spendable — Verify rejects it", weight=True,
            fill=WARN, width=280)

    sc.text(365, 122, 13, "+ witness y", weight=True, fill=INK_AD,
            anchor="middle", width=200)
    sc.line(323, 140, 410, 140, color=MUTED, width=2.5 * sc.u, arrow=True)
    sc.rect(411, 56, 318, 152, r=13, fill=CARD, stroke=BASE, sw=1.6)
    sc.text(431, 86, 12, "2 · ADAPT", weight=True, fill=BASE, spacing=110,
            width=250)
    sc.text(431, 118, 17, "z = ẑ + y", mono=True, width=250)
    sc.text(431, 144, 13.5, "only the witness holder can do this", fill=MUTED,
            width=320)
    sc.rect(431, 158, 278, 36, r=9, fill=REUSED_SOFT, stroke=REUSED, sw=1.2)
    sc.text(445, 181, 13, "now an ordinary signature", weight=True, fill=REUSED,
            width=280)

    sc.text(775, 122, 13, "publish", fill=MUTED, anchor="middle", width=200)
    sc.line(733, 140, 820, 140, color=MUTED, width=2.5 * sc.u, arrow=True)
    sc.rect(821, 56, 318, 152, r=13, fill=WASH, stroke=LINE, sw=1.6)
    sc.text(841, 86, 12, "3 · ON CHAIN", weight=True, fill=MUTED, spacing=110,
            width=250)
    sc.text(841, 118, 17, "σ → chain", mono=True, width=250)
    sc.text(841, 146, 13.5, "it looks like an ordinary payment:", fill=MUTED,
            width=320)
    sc.text(841, 168, 13.5, "no shared hash lock on either chain", fill=MUTED,
            width=340)

    sc.line(980, 208, 980, 236, color=ADAPTOR, width=2.5 * sc.u)
    sc.line(980, 236, 170, 236, color=ADAPTOR, width=2.5 * sc.u)
    sc.line(170, 236, 170, 264, color=ADAPTOR, width=2.5 * sc.u, arrow=True)
    sc.text(600, 228, 13,
            "σ is public — and the counterparty already holds its own σ̂",
            fill=MUTED, anchor="middle", width=700)
    sc.rect(1, 264, 538, 128, r=13, fill=ADAPTOR_SOFT, stroke=ADAPTOR, sw=1.6)
    sc.text(21, 294, 12, "4 · EXTRACT", weight=True, fill=INK_AD, spacing=110,
            width=250)
    sc.text(21, 326, 17, "y = z − ẑ", mono=True, width=250)
    sc.text(21, 352, 13.5, "subtract the pre-signature it already holds;",
            width=460)
    sc.text(21, 374, 13.5, "the recovered y is checked against Y", width=460)
    sc.line(543, 328, 600, 328, color=REUSED, width=2.5 * sc.u, arrow=True)
    sc.rect(601, 264, 538, 128, r=13, fill=REUSED_SOFT, stroke=REUSED, sw=1.6)
    sc.text(621, 294, 12, "THE OTHER LEG", weight=True, fill=REUSED,
            spacing=110, width=250)
    sc.text(621, 326, 17, "settles with the same secret", weight=True, width=460)
    sc.text(621, 352, 13.5, "y adapts the matching pre-signature,", width=460)
    sc.text(621, 374, 13.5, "which is then published on the other chain",
            width=460)

    para_box(s, L, 5900000, W, 420000,
             "[[hb|Both legs settle, or neither does]] — no trusted third "
             "party, and no shared hash lock tying the chains together.",
             size=15.0, bold=True, line=1.25)
    return s


def slide_method(prs, V, meta, idx, total):
    """Slide 4 -- the layered method and the one substitution."""
    s = new_slide(prs, meta, idx, total)
    head(s, "The method · how a basic signature becomes an adaptor signature",
         "One substitution in the challenge hash — [[hi|everything else is reuse]]")

    sc = Scene(s, L, BODY_Y, W, (0, 0, 1160, 300))
    sc.rect(1, 10, 700, 80, r=9, fill=ADAPTOR_SOFT, stroke=ADAPTOR, sw=1.8)
    sc.text(22, 34, 12, "NEW — THE ADAPTOR LAYER", weight=True, fill=ADAPTOR,
            spacing=100, width=400)
    sc.text(22, 60, 18, "PreSign · PreVerify · Adapt · Extract", weight=True,
            width=560)
    sc.text(22, 80, 13,
            "+ byte-level wire encoding, the tests, and the atomic-swap "
            "application", fill=MUTED, width=660)

    sc.rect(1, 102, 700, 70, r=9, fill=BASE_SOFT, stroke=BASE, sw=1.8)
    sc.text(22, 126, 12, "SIMPLIFIED BASE SIGNATURE", weight=True, fill=BASE,
            spacing=100, width=400)
    sc.text(22, 152, 18, "KeyGen · Sign · Verify", weight=True, width=500)

    sc.rect(1, 184, 700, 70, r=9, fill=WASH, stroke=MUTED, sw=1.5)
    sc.text(22, 208, 12, "REUSED UNCHANGED — CRYSTALS-Dilithium reference",
            weight=True, fill=MUTED, spacing=100, width=560)
    sc.text(22, 234, 18, "NTT · polynomial arithmetic · SHAKE · sampling",
            weight=True, width=620)
    sc.text(22, 282, 13.5,
            "zero upstream functions modified — shown by a two-branch diff, "
            "not asserted", fill=MUTED, width=700)

    sc.line(756, 50, 702, 50, color=ADAPTOR, width=2.5 * sc.u, arrow=True)
    sc.rect(757, 10, 402, 244, r=11, fill=ADAPTOR_SOFT, stroke=ADAPTOR, sw=1.6,
            dash=MSO_LINE_DASH_STYLE.DASH)
    sc.text(779, 36, 12, "THE ONE SUBSTITUTION", weight=True, fill=INK_AD,
            spacing=100, width=380)
    sc.text(779, 70, 14.5, "Sign", mono=True, fill=MUTED, width=120)
    sc.text(852, 70, 14.5, "c = H(pk, w, M)", mono=True, width=300)
    sc.text(779, 96, 14.5, "PreSign", mono=True, fill=MUTED, width=120)
    sc.text(852, 96, 14.5, "c = H(pk, w + Y, M)", mono=True, width=300)
    sc.text(852, 118, 12.5, "↑ the statement, folded into the challenge",
            weight=True, fill=INK_AD, width=320)
    sc.line(779, 136, 1137, 136, color=ADAPTOR, width=1.0 * sc.u)
    sc.text(779, 160, 13, "H is instantiated in two steps:", fill=MUTED,
            width=340)
    sc.text(779, 190, 14.5, "c̃ = SHAKE256(pk, w [+ Y], M)", mono=True, width=380)
    sc.text(779, 216, 14.5, "c  = SampleInBall(c̃)", mono=True, width=380)
    sc.text(779, 240, 12.5, "the digest c̃ is what travels on the wire",
            fill=MUTED, width=380)

    cw = (W - GAP) / 2
    cy = 5250000
    _, tf = card(s, L, cy, cw, 1150000, fill=REUSED_SOFT, stroke=REUSED)
    write(tf, "Built a second time, independently", size=13.0, bold=True,
          color=INK, first=True)
    write(tf, "C and Rust agree [[b|byte for byte]] on a pinned known-answer "
              f"digest: [[m|{V['katDigestHead']}…{V['katDigestTail']}]]",
          size=11.5, color=INK, space_before=6, line=1.3)

    caveat(s, L + cw + GAP, cy, cw, 620000,
           "The digest binds packed [[warn|outputs]] over four fixed vectors, "
           "not internals; PreVerify and Extract are asserted per vector, not "
           "hashed.", size=10.5)
    chipline(s, L + cw + GAP, cy + 700000, [
        (f"functional · 1000 iterations × 3 parameter sets", REUSED),
        (f"tamper · every one of {V['tamperFlips']} signature bytes", REUSED),
    ], gap=80000, size=8.5, height=240000)
    return s


def slide_walk(prs, V, meta, idx, total):
    """Slide 5 -- DEMO A, the swap board (all four beats shown at once)."""
    s = new_slide(prs, meta, idx, total)
    head(s, f"Demonstration 1 of 2 · the protocol, stepped on screen · object "
            f"sizes from run {V['stageTwoRunId']}",
         "Step it, [[hi|break it]], and watch the secret fall out")

    unit = (W - 2 * 150000) / 3.15
    ax, aw = L, unit
    bx, bw = L + unit + 150000, unit * 1.15
    cx, cw = bx + unit * 1.15 + 150000, unit
    top, hgt = BODY_Y, 3080000

    # --- Alice ---------------------------------------------------------
    card(s, ax, top, aw, hgt)
    para_box(s, ax + 160000, top + 130000, aw - 320000, 260000,
             "Alice   [[mu|holds the witness]]", size=13.0, bold=True)
    y = top + 440000
    for text, colour, secret in (
        ("sk₁", BASE, False),
        ("witness y", ADAPTOR, True),
        (f"Y · {V['cfgThreeMsgY']} B", BASE, False),
        (f"σ̂₂ from Bob · {V['cfgThreeMsgPreSigB']} B", BASE, False),
        ("σ₂ adapted", BASE, False),
    ):
        chip(s, ax + 160000, y, text, color=INK_AD if secret else INK,
             border=colour, fill=WASH, size=8.5, height=230000)
        y += 280000
    para_box(s, ax + 160000, top + 1950000, aw - 320000, 700000,
             "Knows y, so it completes Bob’s pre-signature and publishes. "
             "[[b|It takes Bob’s coin.]]", size=11.0, color=MUTED, line=1.32)
    para_box(s, ax + 160000, top + 2660000, aw - 320000, 260000,
             "σ₂ ← Adapt((Y, y), pk₂, σ̂₂, m₂)", size=9.5, color=INK_AD,
             mono=True)

    # --- the two chains ------------------------------------------------
    ch = 900000
    for cy, name, verdict, body in (
        (top + 420000, "chain 2", "SETTLED — Alice claimed",
         f"tx₂ + σ₂ · {V['cfgThreeMsgChainB']} B published — pays Alice"),
        (top + 420000 + ch + 200000, "chain 1", "SETTLED — Bob claimed",
         f"tx₁ + σ₁ · {V['cfgThreeMsgChainA']} B published — pays Bob"),
    ):
        card(s, bx, cy, bw, ch, fill=REUSED_SOFT, stroke=REUSED)
        para_box(s, bx + 160000, cy + 140000, bw - 320000, 250000,
                 name, size=11.5, bold=True)
        para_box(s, bx + 160000, cy + 140000, bw - 320000, 250000,
                 verdict, size=8.5, color=REUSED, mono=True,
                 align=PP_ALIGN.RIGHT)
        para_box(s, bx + 160000, cy + 430000, bw - 320000, 400000,
                 body, size=9.0, color=INK, mono=True, line=1.3)
    stamp = rect(s, bx + bw / 2 - 700000, top + 2500000, 1400000, 430000,
                 fill=WARN_SOFT, stroke=WARN, stroke_w=38000, radius=110000,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    stf = stamp.text_frame
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(stf, "REJECTED", size=16.0, color=WARN, bold=True,
          align=PP_ALIGN.CENTER, spacing=50, first=True)
    stamp.rotation = 354.0

    # --- Bob -----------------------------------------------------------
    card(s, cx, top, cw, hgt)
    para_box(s, cx + 160000, top + 130000, cw - 320000, 260000,
             "Bob   [[mu|learns the witness]]", size=13.0, bold=True)
    y = top + 440000
    for text, colour, secret in (
        ("sk₂", BASE, False),
        ("Y", BASE, False),
        (f"π · {V['cfgThreeMsgPi']} B", BASE, False),
        (f"σ̂₁ from Alice · {V['cfgThreeMsgPreSigA']} B", BASE, False),
        ("σ̂₂ it made itself", BASE, False),
        ("y′ extracted", ADAPTOR, True),
    ):
        chip(s, cx + 160000, y, text, color=INK_AD if secret else INK,
             border=colour, fill=WASH, size=8.5, height=230000)
        y += 260000
    para_box(s, cx + 160000, top + 2050000, cw - 320000, 400000,
             "Tries to spend the pre-signature it is holding — [[b|ordinary "
             "Verify refuses it]].", size=10.5, color=MUTED, line=1.3)
    para_box(s, cx + 160000, top + 2440000, cw - 320000, 240000,
             "Verify(σ̂₁, m₁, pk₁) → ⊥", size=9.0, color=INK_AD, mono=True)
    para_box(s, cx + 160000, top + 2690000, cw - 320000, 300000,
             "Needs [[b|nothing further from Alice]]: it reads σ₂ off chain 2 "
             "and adapts σ̂₁ with what comes out.", size=10.0, color=MUTED,
             line=1.25)

    # --- the rail ------------------------------------------------------
    rw = (W - 3 * 130000) / 4
    ry = 5470000
    for k, (label, text, colour) in enumerate((
        ("1 · ABORT GATE",
         "Bob pre-signed only after π and PreVerify(σ̂₁) passed. Neither "
         "pre-signature is spendable by anyone.", REUSED),
        ("2 · THE TRIPWIRE",
         "A pre-signature is not a signature — that is what stops Bob "
         "claiming early.", REUSED),
        ("3 · PUBLISH",
         "On chain it is an ordinary payment: no script, no hash lock, "
         "nothing to see.", REUSED),
        ("4 · LEAK, THEN CLAIM",
         "Ext returns y′; Bob adapts σ̂₁ with that value and publishes. Both "
         "settle, or neither does.", ADAPTOR),
    )):
        rx = L + k * (rw + 130000)
        rect(s, rx, ry, rw, 38000, fill=colour, stroke=None)
        para_box(s, rx, ry + 110000, rw, 230000, label, size=8.5,
                 color=colour, mono=True, bold=True, spacing=90)
        para_box(s, rx, ry + 360000, rw, 620000, text, size=10.0,
                 color=MUTED, line=1.3)
    return s


def slide_time(prs, V, meta, idx, total, fig):
    """Slide 6 -- cost in time."""
    s = new_slide(prs, meta, idx, total)
    head(s, "Cost in time · Simplified Dilithium-III · paired, interleaved "
            "measurements",
         "The adaptor layer is [[hb|nearly free]]. The [[hi|proof]] is not.")

    colw = (W - GAP) / 2
    rx = L + colw + GAP

    sw = (colw - 150000) / 2
    sh = 1150000
    stats = ((f"+{V['ovPreSign']}%", "PreSign over Sign", INK_AD),
             (f"+{V['ovPreVerify']}%", "PreVerify over Verify", INK_AD),
             (f"+{V['ovAdapt']}%", "Adapt over Verify", INK_AD),
             ("cheapest", "Extract — no basic analogue", REUSED))
    for k, (value, label, colour) in enumerate(stats):
        x = L + (k % 2) * (sw + 150000)
        y = BODY_Y + (k // 2) * (sh + 150000)
        _, tf = card(s, x, y, sw, sh)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(tf, value, size=26.0, color=colour, bold=True, first=True,
              line=1.05)
        write(tf, label, size=11.0, color=MUTED, space_before=5)
    caveat(s, L, BODY_Y + 2 * sh + 320000, colw, 900000,
           "Rejection attempts are [[warn|counted]], never inferred from "
           f"timing: {V['rejAttBase']} per Sign against {V['rejAttLas']} per "
           "PreSign, and the run aborts unless every timed call matches theory.",
           size=10.5)

    fh = 3150000
    rect(s, rx, BODY_Y, colw, fh, fill=CARD, stroke=LINE, radius=110000,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    picture_fit(s, fig, rx + 110000, BODY_Y + 110000, colw - 220000, fh - 220000)
    caveat(s, rx, BODY_Y + fh + 200000, colw, 700000,
           f"In a whole swap the [[warn|proof]] takes "
           f"[[warn|{V['cfgThreeProofPct']}%]] of end-to-end time — the "
           "signature is not the bottleneck.", size=10.5)
    return s


def slide_bytes(prs, V, meta, idx, total):
    """Slide 7 -- cost in bytes."""
    s = new_slide(prs, meta, idx, total)
    head(s, "Result 2 · communication · the price of post-quantum",
         "The cost is not adaptor computation — it is [[hi|bytes on the wire]]")

    # legend
    lx = L
    for swatch, label in ((BASE, "classical ECDSA adaptor"),
                          (ADAPTOR, "LAS")):
        rect(s, lx, BODY_Y + 40000, 130000, 130000, fill=swatch, stroke=None,
             radius=30000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        para_box(s, lx + 190000, BODY_Y, 2600000, 220000, label,
                 size=10.0, color=MUTED, wrap=False)
        lx += 190000 + len(label) * 10.0 * 0.58 * 12700 + 260000
    para_box(s, lx, BODY_Y, 4000000, 220000,
             "one shared byte scale · Simplified Dilithium-II",
             size=10.0, color=MUTED)

    labw, trkx, trkw = 1450000, 2100000, 6000000
    bars = (
        ("classical\nsignature", V["wClSig"], BASE, f"{V['clSigBytes']} B", ""),
        ("LAS\nsignature", V["wLasSig"], ADAPTOR, f"{V['sigBytesTwo']} B",
         f"≈ {V['clRatioSig']}× classical"),
        ("classical\npublic key", V["wClPk"], BASE, f"{V['clPkBytes']} B", ""),
        ("LAS\npublic key", V["wLasPk"], ADAPTOR, f"{V['pkBytesTwo']} B",
         f"≈ {V['clRatioPk']}× classical"),
    )
    by = 2560000
    for label, width_pct, colour, value, note in bars:
        para_box(s, L, by - 20000, labw, 320000, label, size=10.0,
                 color=MUTED, align=PP_ALIGN.RIGHT, line=1.2)
        rect(s, trkx, by, trkw, 300000, fill=WASH, stroke=LINE, radius=55000,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        w = max(40000, trkw * float(width_pct) / 100.0)
        rect(s, trkx, by, w, 300000, fill=colour, stroke=None, radius=45000,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt = f"[[m|{value}]]" + (f"  [[mu|{note}]]" if note else "")
        para_box(s, trkx + trkw + 130000, by + 30000,
                 R - (trkx + trkw + 130000), 300000, txt, size=10.5,
                 color=INK, anchor=MSO_ANCHOR.TOP)
        by += 380000

    cw = (W - GAP) / 2
    cy = 4230000
    sp, tf = card(s, L, cy, cw, 1420000)
    write(tf, f"Inside a LAS signature ({V['sigBytesTarget']} B at the target "
              "set)", size=12.0, bold=True, color=INK, first=True)
    zp = float(V["zPctTarget"])
    sw = (cw - 320000)
    rect(s, L + 160000, cy + 480000, sw * zp / 100.0, 300000, fill=ADAPTOR,
         stroke=None, radius=50000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, L + 160000 + sw * zp / 100.0 + 20000, cy + 480000,
         max(60000, sw * (100 - zp) / 100.0), 300000, fill=BASE, stroke=None,
         radius=50000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    para_box(s, L + 160000, cy + 850000, sw, 480000,
             f"[[hi|response z — {V['zPctTarget']}%]]   ·   "
             f"[[hb|challenge digest c̃ — {V['zRestPct']}%]]\n"
             "The challenge travels only as its digest, so [[b|optimising LAS "
             "means optimising z]].", size=10.5, color=MUTED, line=1.3)

    sp, tf = card(s, L + cw + GAP, cy, cw, 1420000, fill=ADAPTOR_SOFT,
                  stroke=ADAPTOR)
    write(tf, "And an adaptor swap adds one public object", size=12.0,
          bold=True, color=INK, first=True)
    write(tf, f"The statement [[b|Y]] — {V['stmtBytesTarget']} B, exactly the "
              "size of a public key — plus a signature-sized pre-signature.",
          size=11.0, color=INK, space_before=6, line=1.3)
    write(tf, "Those are the only objects an adaptor scheme sends that a basic "
              "signature does not.", size=11.0, color=MUTED, space_before=5,
          line=1.3)

    caveat(s, L, 5810000, W, 600000,
           "[[warn|Fairness:]] the classical baseline is a "
           "functionality-matched ECDSA adaptor signature (not plain ECDSA), "
           "compared at Simplified Dilithium-II — an engineering match to "
           "≈128-bit, [[warn|not]] a formal security-equivalence claim.",
           size=10.5)
    return s


def slide_bitcoin(prs, V, meta, idx, total):
    """Slide 8 -- DEMO B, the node differential."""
    s = new_slide(prs, meta, idx, total)
    head(s, f"Demonstration 2 of 2 · a real Bitcoin Core {V['btcNodeTag']} "
            "node, on regtest",
         "It [[hb|fits]] the size limits. It fails the [[hi|verification "
         "rules]] — so I changed one.")

    lw = (W - GAP) / 2.35
    rw = W - GAP - lw
    rx = L + lw + GAP

    _, tf = card(s, L, BODY_Y, lw, 1500000, fill=BASE_SOFT, stroke=BASE)
    write(tf, "1 · Carriage, stock node", size=12.5, bold=True, color=INK,
          first=True)
    write(tf, f"A lattice-carrying spend: [[b|{V['btcMeasCarriageVsize']} vB]], "
              f"refused by default relay policy "
              f"([[m|{V['btcMeasCarriagePolicy']}]]) — yet consensus-valid, "
              "and [[b|mined]].", size=10.5, color=INK, space_before=6, line=1.3)
    write(tf, "Standardness and validity are different questions; only the "
              "second decides feasibility.", size=10.0, color=MUTED,
          space_before=5, line=1.3)

    _, tf = card(s, L, BODY_Y + 1620000, lw, 1900000, fill=REUSED_SOFT,
                 stroke=REUSED)
    write(tf, "2 · Verification, patched node", size=12.5, bold=True,
          color=INK, first=True)
    write(tf, "One of the reserved [[m|OP_SUCCESS]] opcodes, "
              f"[[m|{V['btcLasOpcode']}]], defined as lattice verification — "
              "so a spend can be authorised by a LAS signature alone.",
          size=10.5, color=INK, space_before=6, line=1.3)
    write(tf, "A [[b|whole two-leg swap]] settles on it: the adapted signature "
              "comes back out of the [[b|mined]] transaction, and "
              "[[b|Extract]] recovers the secret. Verdicts are "
              "[[b|consensus]], not policy.", size=10.0, color=MUTED,
          space_before=5, line=1.3)

    caveat(s, L, 5880000, lw, 550000,
           "[[warn|Honest limits:]] a patched node is not Bitcoin, so “cannot "
           "settle on Bitcoin as it stands” still holds; and the security of "
           "the new rule is [[warn|not]] analysed.", size=9.5)

    # the differential
    _, tf = card(s, rx, BODY_Y, rw, 760000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, f"the [[b|same spend]], byte for byte — [[b|{V['btcLasVsize']} "
              f"vB]], {V['btcLasItems']} witness items, [[b|no elliptic-curve "
              "signature at all]] — put to two clients of the same release",
          size=10.5, color=MUTED, align=PP_ALIGN.CENTER, line=1.3, first=True)
    forky = BODY_Y + 760000
    line(s, rx + rw / 2, forky, rx + rw / 2, forky + 130000, color=LINE,
         width=19050)
    line(s, rx + rw / 4, forky + 130000, rx + 3 * rw / 4, forky + 130000,
         color=LINE, width=19050)
    line(s, rx + rw / 4, forky + 130000, rx + rw / 4, forky + 300000,
         color=LINE, width=19050)
    line(s, rx + 3 * rw / 4, forky + 130000, rx + 3 * rw / 4, forky + 300000,
         color=LINE, width=19050)

    nw = (rw - 150000) / 2
    ny = forky + 380000
    nh = 2100000
    for k, (title, sub, rows, border) in enumerate((
        (f"patched Core {V['btcNodeTag']}",
         f"{V['btcLasOpcode']} verifies the LAS signature",
         (("ACCEPTED · MINED", "the valid spend", REUSED, REUSED_SOFT),
          (f"REJECTED ×{V['btcLasControls']}", "every negative control", WARN,
           WARN_SOFT)), ADAPTOR),
        (f"stock Core {V['btcNodeTag']}",
         f"{V['btcLasOpcode']} is still OP_SUCCESS",
         (("ACCEPTED", "the valid spend — same block", REUSED, REUSED_SOFT),
          (f"ACCEPTED ×{V['btcLasControls']}", "every negative control",
           REUSED, REUSED_SOFT)), LINE),
    )):
        nx = rx + k * (nw + 150000)
        card(s, nx, ny, nw, nh, stroke=border)
        para_box(s, nx + 160000, ny + 130000, nw - 320000, 250000, title,
                 size=11.5, bold=True)
        para_box(s, nx + 160000, ny + 390000, nw - 320000, 240000, sub,
                 size=9.5, color=MUTED, mono=True)
        vy = ny + 700000
        for key, text, colour, bg in rows:
            rect(s, nx + 160000, vy, nw - 320000, 620000, fill=bg, stroke=None,
                 radius=80000, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            para_box(s, nx + 260000, vy + 90000, nw - 520000, 230000, key,
                     size=9.0, color=colour, mono=True, bold=True, spacing=50)
            para_box(s, nx + 260000, vy + 320000, nw - 520000, 240000, text,
                     size=10.0, color=INK)
            vy += 700000

    para_box(s, rx, ny + nh + 180000, rw, 700000,
             "Both clients accept the honest spend. Only the patched one can "
             "tell a bad LAS signature from a good one — so [[b|the difference "
             "between the columns is the new rule, and nothing else]].",
             size=10.5, color=MUTED, align=PP_ALIGN.CENTER, line=1.3)
    return s


def slide_wrong(prs, V, meta, idx, total, fig):
    """Slide 9 -- what measurement overturned."""
    s = new_slide(prs, meta, idx, total)
    head(s, "Corrected by measurement, not by argument",
         "Two beliefs [[hi|measurement overturned]] — and three suggestions I "
         "closed instead")

    cw = (W - GAP) / 2
    rx = L + cw + GAP

    fh = 2050000
    rect(s, L, BODY_Y, cw, fh, fill=CARD, stroke=LINE, radius=110000,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    picture_fit(s, fig, L + 110000, BODY_Y + 110000, cw - 220000, fh - 220000)

    _, tf = card(s, L, BODY_Y + fh + 150000, cw, 1330000)
    write(tf, "“That can be future work”", size=12.5, bold=True, color=INK,
          first=True)
    write(tf, "Three optimisations I could have listed as suggestions, run and "
              "closed instead:", size=10.5, color=INK, space_before=5, line=1.3)
    cy = BODY_Y + fh + 150000 + 700000
    chipline(s, L + 160000, cy, [
        ("statement compression — fails at Extract, every depth tested", INK_AD),
    ], size=8.0, height=230000)
    chipline(s, L + 160000, cy + 270000, [
        ("batching — amortises the wrong cost", INK_AD),
        ("optimistic verifier — fraud proof cannot be mined", INK_AD),
    ], gap=80000, size=8.0, height=230000)

    _, tf = card(s, rx, BODY_Y, cw, 1000000, fill=REUSED_SOFT, stroke=REUSED)
    write(tf, "“It can never fit on chain”", size=12.5, bold=True, color=INK,
          first=True)
    write(tf, "I called the cause structural and inferred optimisation could "
              "not close it. [[warn|Wrong]] — a real client mined full "
              f"verification at [[b|{V['gasOptCapPct']}%]] of the cap.",
          size=11.0, color=INK, space_before=6, line=1.3)

    c2y = BODY_Y + 1100000
    c2h = 2380000
    _, tf = card(s, rx, c2y, cw, c2h, fill=REUSED_SOFT, stroke=REUSED)
    write(tf, "“The base scheme must be simplified”", size=12.5, bold=True,
          color=INK, first=True)
    write(tf, "Built a third time on the [[b|full FIPS 204 ML-DSA path]]: "
              "PreSign and PreVerify are new algorithms, but [[b|Verify is "
              f"not]] — the stock verifier accepts "
              f"[[b|{V['mldsaRepairedP']}/{V['mldsaIters']}]] adapted "
              "signatures.", size=10.5, color=INK, space_before=6, line=1.3)

    ty = c2y + 1080000
    tw = cw - 320000
    cols = (tw * 0.46, tw * 0.27, tw * 0.27)
    rows = (("", "simplified", "on ML-DSA"),
            ("attempts per PreSign", V["mldsaAttemptsSimp"], V["mldsaAttempts"]),
            ("signature", f"{V['mldsaSigBytesSimp']} B",
             f"{V['mldsaSigBytes']} B"),
            ("statement Y", f"{V['mldsaStmtBytes']} B",
             f"{V['mldsaStmtBytes']} B"))
    for r, row in enumerate(rows):
        yy = ty + r * 240000
        if r:
            line(s, rx + 160000, yy, rx + 160000 + tw, yy, color=LINE,
                 width=9525)
        xx = rx + 160000
        for c, cell in enumerate(row):
            head_row = r == 0
            para_box(s, xx, yy + 40000, cols[c], 220000, cell,
                     size=8.5 if head_row else 9.0,
                     color=MUTED if (head_row or c == 0) else INK,
                     mono=not head_row and c > 0,
                     caps=head_row, spacing=90 if head_row else None,
                     align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
            xx += cols[c]
    para_box(s, rx + 160000, ty + 4 * 240000 + 30000, tw, 240000,
             "[[b|Y]] is unchanged, so it is the [[b|dominant remaining size "
             "target]].", size=10.5, color=INK)

    caveat(s, L, 5830000, W, 570000,
           "On-chain gas uses [[warn|one fixed signature instance]] per "
           "measured parameter set, so instance-to-instance variation is "
           "unquantified; Dilithium-V is derived to exceed one transaction, "
           "never measured there; and the ML-DSA build is a [[warn|functional]] "
           "demonstration — its unforgeability is not analysed.", size=10.0)
    return s


def slide_answers(prs, V, meta, idx, total):
    """Slide 10 -- back to the opening questions."""
    s = new_slide(prs, meta, idx, total)
    head(s, "Back to the opening questions",
         "Post-quantum atomic swaps are [[hb|buildable today]] — the bill "
         "arrives in [[hi|bytes]]")

    cw = (W - 2 * GAP) / 3
    cy, ch = BODY_Y, 2900000
    cards = (
        ("YES", REUSED, REUSED_SOFT, "Can we build one?",
         "[[b|Yes]] — and twice.",
         "Reuse of standardised primitives with zero upstream functions "
         "modified, reproduced independently in a second language to a pinned "
         "byte-identical digest."),
        ("BYTES", ADAPTOR, ADAPTOR_SOFT, "What does it cost?",
         f"Almost nothing in computation: at most [[b|+{V['ovMaxAll']}%]] for "
         "the adaptor layer.",
         "The bill arrives in [[b|bytes]] — and in a whole swap it is the "
         "[[b|proof]], not the signature, that dominates the time."),
        ("SPLIT", BASE, BASE_SOFT, "Will a chain take it?",
         "[[b|Ethereum: yes]] — a real client mined full verification inside "
         "one transaction.",
         "[[b|Bitcoin: not as it stands]] — carriage works on a stock node, "
         "but verification needs a new consensus rule; I patched one, and "
         "settled a whole two-leg swap on it."),
    )
    for k, (badge, colour, soft, question, lead, body) in enumerate(cards):
        x = L + k * (cw + GAP)
        card(s, x, cy, cw, ch, fill=soft, stroke=colour)
        chip(s, x + 160000, cy + 150000, badge, color=colour, border=colour,
             fill=CARD, size=9.0, height=260000)
        para_box(s, x + 160000, cy + 490000, cw - 320000, 320000, question,
                 size=13.0, bold=True)
        para_box(s, x + 160000, cy + 860000, cw - 320000, 620000, lead,
                 size=11.5, color=INK, line=1.32)
        para_box(s, x + 160000, cy + 1520000, cw - 320000, 1230000, body,
                 size=11.0, color=MUTED, line=1.32)

    para_box(s, L, 5350000, W, 900000,
             "Still open, specifically: a relation whose [[hi|statement is "
             "small by construction]]; a proof system succinct at [[hi|this]] "
             "scale; and the [[hi|security]] of the ML-DSA variant and of the "
             "consensus rule.", size=14.0, bold=True, line=1.3)
    return s


# ================================================================== main ====
def build(output: pathlib.Path) -> Presentation:
    values = gen_slides.read_macros()
    values.update(gen_slides.derive(values))
    meta = deck_meta(values)
    if len(meta) != 10:
        sys.exit(f"expected 10 slides in the template, found {len(meta)}")

    prs = Presentation(str(UOM_PPTX))
    for sld in list(prs.slides._sldIdLst):          # drop the template samples
        prs.part.drop_rel(sld.get(qn("r:id")))
        prs.slides._sldIdLst.remove(sld)

    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = pathlib.Path(tmp)
        fig_timing = figure_png("fig_timing", tmpdir)
        fig_onchain = figure_png("fig_onchain", tmpdir)
        n = len(meta)
        slide_title(prs, values, meta[0], 1, n)
        slide_why(prs, values, meta[1], 2, n)
        slide_mechanism(prs, values, meta[2], 3, n)
        slide_method(prs, values, meta[3], 4, n)
        slide_walk(prs, values, meta[4], 5, n)
        slide_time(prs, values, meta[5], 6, n, fig_timing)
        slide_bytes(prs, values, meta[6], 7, n)
        slide_bitcoin(prs, values, meta[7], 8, n)
        slide_wrong(prs, values, meta[8], 9, n, fig_onchain)
        slide_answers(prs, values, meta[9], 10, n)
        prs.save(str(output))
    return prs


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("-o", "--output", default=str(OUTPUT))
    args = ap.parse_args()
    out = pathlib.Path(args.output)
    build(out)
    size = out.stat().st_size / 1024
    print(f"wrote {out.relative_to(ROOT)} (10 slides, {size:.0f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
