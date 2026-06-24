# -*- coding: utf-8 -*-
"""
Build the LAS-on-Dilithium EVALUATION — DEEP ANALYSIS deck (.pptx).

A companion to the design-defence deck. It presents *all* evaluation metrics of
the project and analyses them: signature-level timings, rejection/acceptance,
object sizes, the 2x2 cross-scheme matrix (vs Dilithium-3 and vs a classical
adaptor), application-level payloads, AMHL-vs-K, and on-chain gas (incl. the
native-verification cost experiment), plus correctness/robustness — with the
three headline findings and an honest threats-to-validity pass.

All numbers are sourced from docs/LAS.md Section 8 (representative runs; the
ratios are the result). Reproducible: re-run to regenerate the .pptx.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------- palette
BG      = RGBColor(0x0D, 0x10, 0x30)
PANEL   = RGBColor(0x16, 0x1B, 0x45)
CARD    = RGBColor(0x1E, 0x24, 0x57)
CARD2   = RGBColor(0x26, 0x2D, 0x66)
ACCENT  = RGBColor(0x49, 0xE3, 0xCE)   # cyan/mint -> LAS / positive
AMBER   = RGBColor(0xFF, 0xB4, 0x54)   # amber     -> adaptor ops / highlight
CORAL   = RGBColor(0xFF, 0x6B, 0x6B)   # coral     -> cost / classical-exotic / caution
VIOLET  = RGBColor(0x9B, 0x8C, 0xFF)   # violet    -> Dilithium-3
TEXT    = RGBColor(0xEC, 0xEE, 0xFF)
MUTED   = RGBColor(0x9A, 0xA0, 0xD0)
DIM     = RGBColor(0x70, 0x78, 0xB0)
LINE    = RGBColor(0x2E, 0x35, 0x6F)
CODEBG  = RGBColor(0x0A, 0x0D, 0x28)
LATTICE = RGBColor(0x29, 0x31, 0x68)

HEAD = "Georgia"
BODY = "Calibri"
MONO = "Consolas"
CENTER = PP_ALIGN.CENTER
RR = MSO_SHAPE.ROUNDED_RECTANGLE
EMU_IN = 914400
SW, SH = 13.333, 7.5

# ----------------------------------------------------------------------------- helpers
def set_bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color

def no_shadow(sh): sh.shadow.inherit = False

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(sh)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    if radius is not None and shape == RR:
        try: sh.adjustments[0] = radius
        except Exception: pass
    return sh

def _set_margins(tf, m=0.06):
    tf.margin_left = Inches(m); tf.margin_right = Inches(m)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)

def _run(p, text, size, color, font=BODY, bold=False, italic=False, spc=None):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.name = font; f.bold = bold; f.italic = italic
    f.color.rgb = color
    if spc is not None:
        rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(int(spc * 100)))
    return r

def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True, m=0.06):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    _set_margins(tf, m); tf.paragraphs[0].alignment = align
    return tb, tf

def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    if line is not None: p.line_spacing = line
    return p

def simple(s, x, y, w, h, text, size, color, font=BODY, bold=False, italic=False,
           align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line=None, m=0.06):
    tb, tf = textbox(s, x, y, w, h, anchor=anchor, align=align, m=m)
    p = para(tf, first=True, align=align, space_after=0, line=line)
    _run(p, text, size, color, font=font, bold=bold, italic=italic, spc=spc)
    return tb

def bullets(s, x, y, w, h, items, size=14.5, gap=7, color=TEXT, marker_color=ACCENT,
            marker="▸", font=BODY, line=1.0, anchor=MSO_ANCHOR.TOP):
    tb, tf = textbox(s, x, y, w, h, anchor=anchor)
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0), space_after=gap, line=line)
        if marker: _run(p, marker + "  ", size, marker_color, font=font, bold=True)
        if isinstance(it, str): _run(p, it, size, color, font=font)
        else:
            for (t, o) in it:
                _run(p, t, o.get("size", size), o.get("color", color), font=o.get("font", font),
                     bold=o.get("bold", False), italic=o.get("italic", False))
    return tb

def dot_grid(s, x0, y0, cols, rows, step=0.30, r=0.040, color=LATTICE):
    for i in range(cols):
        for j in range(rows):
            o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x0+i*step), Inches(y0+j*step), Inches(r), Inches(r))
            no_shadow(o); o.line.fill.background(); o.fill.solid(); o.fill.fore_color.rgb = color

def kicker_title(s, kicker, title, accent=ACCENT, title_size=30, ty=0.74, tw=12.1):
    simple(s, 0.62, 0.42, 11.0, 0.32, kicker.upper(), 12.5, accent, font=BODY, bold=True, spc=2.4)
    simple(s, 0.6, ty, tw, 0.95, title, title_size, TEXT, font=HEAD, bold=True)

def footer(s, n):
    simple(s, 0.62, 7.06, 9.5, 0.3, "LAS on CRYSTALS-Dilithium  ·  Evaluation — Deep Analysis  ·  docs/LAS.md §8",
           9, DIM, font=BODY)
    simple(s, 12.0, 7.06, 0.9, 0.3, f"{n:02d}", 9.5, MUTED, font=MONO, bold=True, align=PP_ALIGN.RIGHT)

def base(prs, kicker, title, accent=ACCENT, n=0, motif=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    if motif: dot_grid(s, 12.05, 0.30, 4, 3)
    kicker_title(s, kicker, title, accent=accent)
    footer(s, n)
    return s

def notes(s, txt): s.notes_slide.notes_text_frame.text = txt

def code_panel(s, x, y, w, h, lines, size=12.5, title=None, accent=ACCENT, fill=CODEBG, border=LINE):
    rect(s, x, y, w, h, fill=fill, line=border, lw=1.0, shape=RR, radius=0.04)
    pad = 0.16; yy = y + pad
    if title:
        simple(s, x + pad, yy, w - 2 * pad, 0.3, title, 11, accent, font=BODY, bold=True, spc=1.5); yy += 0.34
    tb, tf = textbox(s, x + pad, yy, w - 2 * pad, h - (yy - y) - pad)
    for i, ln in enumerate(lines):
        p = para(tf, first=(i == 0), space_after=3, line=1.04)
        if isinstance(ln, str): _run(p, ln if ln else " ", size, TEXT, font=MONO)
        else:
            for (t, o) in ln:
                _run(p, t, o.get("size", size), o.get("color", TEXT), font=o.get("font", MONO),
                     bold=o.get("bold", False), italic=o.get("italic", False))
    return tb

def card(s, x, y, w, h, fill=CARD, line=None, radius=0.045):
    return rect(s, x, y, w, h, fill=fill, line=line, lw=1.0, shape=RR, radius=radius)

def statcard(s, x, y, w, h, big, label, col=ACCENT, fill=CARD2, big_size=34, line=None):
    card(s, x, y, w, h, fill=fill, line=(line or col))
    simple(s, x+0.22, y+0.14, w-0.44, h*0.5, big, big_size, col, font=HEAD, bold=True)
    simple(s, x+0.22, y+h*0.5+0.06, w-0.44, h*0.5-0.1, label, 12, TEXT, font=BODY, line=1.04)

# horizontal bar (track + fill), optional end-label
def hbar(s, x, y, track_w, h, value, maxval, color, track=PANEL, valstr=None, val_color=None):
    rect(s, x, y, track_w, h, fill=track, shape=RR, radius=0.5)
    ww = max(0.07, track_w * (value / float(maxval)))
    rect(s, x, y, ww, h, fill=color, shape=RR, radius=0.5)
    if valstr:
        simple(s, x+ww+0.1, y-0.02, 1.7, h+0.04, valstr, 11.5, val_color or color,
               font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# vertical bar growing upward from base_y
def vbar(s, cx, base_y, w, value, maxval, max_h, color, toplabel=None, botlabel=None,
         tl_color=None, label_w=1.2):
    hh = max(0.05, max_h * (value / float(maxval)))
    rect(s, cx, base_y - hh, w, hh, fill=color, shape=RR, radius=0.12)
    if toplabel:
        simple(s, cx + w/2 - label_w/2, base_y - hh - 0.34, label_w, 0.3, toplabel, 12.5,
               tl_color or color, font=MONO, bold=True, align=CENTER)
    if botlabel:
        simple(s, cx + w/2 - label_w/2, base_y + 0.06, label_w, 0.55, botlabel, 10.5, MUTED,
               font=BODY, align=CENTER, line=1.0)

# ============================================================================= build
prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN)); prs.slide_height = Emu(int(SH * EMU_IN))
prs.author = "Royce Steven"; prs.title = "LAS on Dilithium - Evaluation Deep Analysis"

# ----------------------------------------------------------------- 1 | TITLE
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG)
rect(s, 0.0, 0.0, 0.18, SH, fill=ACCENT)
dot_grid(s, 9.2, 0.5, 11, 6, step=0.36, r=0.05, color=LATTICE)
simple(s, 0.7, 1.35, 11.5, 0.4, "MSc CYBERSECURITY THESIS · EVALUATION CHAPTER",
       13, ACCENT, font=BODY, bold=True, spc=2.6)
simple(s, 0.68, 1.95, 12.0, 1.5, "Evaluation: a deep analysis", 50, TEXT, font=HEAD, bold=True)
simple(s, 0.7, 2.95, 12.0, 1.0, "Every metric, measured — and what it actually means", 25, MUTED, font=HEAD, italic=True)
rect(s, 0.72, 4.2, 6.6, 0.022, fill=LINE)
bullets(s, 0.72, 4.45, 11.8, 1.6, [
    [("Two benchmark types", {"bold":True,"color":TEXT,"size":15}),
     ("  (signature + application)  ×  ", {"color":MUTED,"size":15}),
     ("two baselines", {"bold":True,"color":TEXT,"size":15}),
     ("  (vs Dilithium-3, vs a classical adaptor)", {"color":MUTED,"size":15})],
    [("Timings · acceptance · sizes · payloads · on-chain gas · correctness", {"color":ACCENT,"size":15,"bold":True})],
    [("All figures from ", {"color":MUTED,"size":14}),
     ("docs/LAS.md §8", {"color":ACCENT,"size":14,"bold":True}),
     ("  —  absolute values are machine-dependent; ", {"color":MUTED,"size":14}),
     ("the ratios are the result.", {"color":TEXT,"size":14,"bold":True})],
], size=15, gap=9, marker="—", marker_color=AMBER)
simple(s, 0.72, 6.55, 11.0, 0.5, "Royce Steven   ·   Supervisor: Wang Zhipeng   ·   2026", 13, DIM, font=BODY)
notes(s, "This is the evaluation deep-dive. Wang asked for two benchmark types — the signature itself and "
         "the application — each against two baselines: pure Dilithium-3 and a classical adaptor signature. "
         "I present every metric and, crucially, analyse it: what it measures, how it was measured, and what "
         "it means. The framing rule throughout is 'let the data speak' — absolute microsecond numbers are "
         "machine-dependent, so the ratios and the structure are the real findings.")

# ----------------------------------------------------------------- 2 | EVALUATION MAP
s = base(prs, "What was measured", "The evaluation map: metrics × baselines", n=2)
groups = [
    ("SIGNATURE-LEVEL", ACCENT, [
        "per-op timing (8 ops)", "rejection / acceptance rate", "object sizes (3 distinct numbers)"]),
    ("APPLICATION-LEVEL", AMBER, [
        "atomic-swap payload", "AMHL cost vs path length K", "settlement footprint"]),
    ("ON-CHAIN (EVM)", CORAL, [
        "real Solidity swap gas", "claim / fund / refund", "native-verify cost probe"]),
    ("CORRECTNESS", VIOLET, [
        "1000 iters × modes 2/3/5", "4672 byte-flip forgery test", "pinned KAT digest"]),
]
gw, gh = 2.92, 3.05; gx0, gy = 0.62, 2.05; gap = 0.13
for i,(t,col,its) in enumerate(groups):
    cx = gx0 + i*(gw+gap)
    card(s, cx, gy, gw, gh, fill=CARD, line=LINE)
    rect(s, cx, gy, gw, 0.62, fill=PANEL, shape=RR, radius=0.06)
    rect(s, cx, gy+0.40, gw, 0.22, fill=PANEL)   # square off header bottom
    rect(s, cx, gy, 0.10, gh, fill=col)
    simple(s, cx+0.26, gy, gw-0.4, 0.62, t, 12.5, col, font=BODY, bold=True, spc=0.6, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, cx+0.24, gy+0.8, gw-0.42, gh-0.95, its, size=12.5, gap=11, marker="·", marker_color=col, color=TEXT)
# baseline strip
card(s, 0.62, 5.45, 12.0, 1.05, fill=PANEL, line=LINE)
simple(s, 0.85, 5.56, 4.0, 0.3, "TWO BASELINES (Meeting-2 B2)", 11.5, MUTED, font=BODY, bold=True, spc=1.0)
bullets(s, 0.85, 5.9, 11.6, 0.55, [
    [("(i) LAS vs pure Dilithium-3", {"color":VIOLET,"bold":True}),
     ("  — feasibility + adaptor overhead.    ", {"color":TEXT}),
     ("(ii) LAS vs classical ECDSA-adaptor", {"color":CORAL,"bold":True}),
     ("  — the 'price of post-quantum', same machine & compiler.", {"color":TEXT})],
], size=13.5, gap=4, marker="")
notes(s, "The structure of the evaluation. Four metric families: signature-level (timings, acceptance, "
         "sizes), application-level (swap payload and the multi-hop AMHL cost as a function of path length), "
         "on-chain gas on a real EVM, and correctness/robustness. Each is measured against the two baselines "
         "the supervisor required. The classical adaptor — an ECDSA-adaptor from libsecp256k1-zkp — is "
         "benchmarked on the same machine and compiler, so the comparison needs no hardware caveats.")

# ----------------------------------------------------------------- 3 | METHODOLOGY & RIGOR
s = base(prs, "How it was measured", "Methodology & rigour — and where to be careful", n=3)
left = [
    [("Wall-clock µs/op, ", {"color":TEXT}), ("mode 3, 2000 iters, −O3", {"color":ACCENT,"bold":True}),
     ("; medians reported, run-to-run variation a few %.", {"color":TEXT})],
    [("Acceptance measured ", {"color":TEXT}), ("directly", {"color":ACCENT,"bold":True}),
     (" via the ", {"color":TEXT}), ("las_attempts", {"color":ACCENT,"font":MONO}),
     (" counter — not estimated from a timing ratio.", {"color":TEXT})],
    [("Sizes are the ", {"color":TEXT}), ("actual serialised bytes", {"color":ACCENT,"bold":True}),
     (" from serialize.c, validated by test_serde — not formulas.", {"color":TEXT})],
    [("EVM gas is ", {"color":TEXT}), ("deterministic", {"color":ACCENT,"bold":True}),
     (" (not machine-dependent); same harness mirrors the LAS op-set 1:1.", {"color":TEXT})],
]
bullets(s, 0.62, 2.0, 6.2, 3.4, left, size=14, gap=14)
# caveats card
card(s, 7.0, 1.95, 5.62, 4.55, fill=PANEL, line=CORAL)
simple(s, 7.22, 2.1, 5.2, 0.4, "THREATS TO VALIDITY (stated up front)", 12, CORAL, font=BODY, bold=True, spc=0.8)
bullets(s, 7.22, 2.6, 5.25, 3.7, [
    [("libsecp256k1 is constant-time, optimised production code", {"color":TEXT,"bold":True}),
     (" — the timing comparison flatters the classical side.", {"color":MUTED})],
    [("Not a same-security comparison", {"color":TEXT,"bold":True}),
     (": LAS sits at q≈2²³ with smaller dimensions than Dilithium-3.", {"color":MUTED})],
    [("App-level uses a simulated ledger", {"color":TEXT,"bold":True}),
     (" — byte payloads are a proxy; only §8.4 is real EVM gas.", {"color":MUTED})],
    [("The native-verify probe is gas-faithful, not numerically correct", {"color":TEXT,"bold":True}),
     (" — a lower bound (excludes A′ expansion + unpacking).", {"color":MUTED})],
], size=12.5, gap=11, marker="!", marker_color=CORAL)
notes(s, "Before any numbers, the methodology and its limits. Timings are medians over 2000 iterations at "
         "-O3; the rejection rate is measured with a direct counter rather than inferred; sizes are real "
         "serialised bytes; EVM gas is deterministic. The threats-to-validity panel is deliberately on its "
         "own slide and up front: the classical library is hand-optimised constant-time code so it flatters "
         "the classical timings, this is not a same-security comparison, the app-level figures are a "
         "simulated-ledger proxy, and the native-verify number is a gas-faithful lower bound. Naming these "
         "first is what makes the rest credible.")

# ----------------------------------------------------------------- 4 | PER-OP TIMINGS
s = base(prs, "Signature-level · timings", "Per-operation cost — and the adaptor overhead", n=4)
simple(s, 0.62, 1.72, 7.6, 0.35, "Standalone LAS, mode 3 (bench_las3) — µs/op", 12.5, MUTED, font=BODY, italic=True)
ops = [  # (name, value, color)
    ("Sign",      804, ACCENT),
    ("PreSign",   828, AMBER),
    ("Verify",    191, ACCENT),
    ("PreVerify", 197, AMBER),
    ("Adapt",     203, AMBER),
    ("Ext",        68, ACCENT),
    ("KeyGen",     78, ACCENT),
    ("Setup",      58, ACCENT),
]
bx, by0, bh, bgap, track = 2.45, 2.2, 0.33, 0.135, 4.9
maxv = 900
for i,(nm,v,col) in enumerate(ops):
    yy = by0 + i*(bh+bgap)
    simple(s, 0.62, yy-0.02, 1.75, bh+0.04, nm, 12.5, col, font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    hbar(s, bx, yy, track, bh, v, maxv, col, valstr=f"{v}", val_color=col)
# pairing brackets / callout
card(s, 8.05, 2.0, 4.55, 2.05, fill=CARD2, line=ACCENT)
simple(s, 8.27, 2.14, 4.1, 0.4, "THE ADAPTOR OVERHEAD ≈ 0", 12, ACCENT, font=BODY, bold=True, spc=0.8)
bullets(s, 8.27, 2.58, 4.15, 1.4, [
    [("PreSign ≈ Sign", {"color":AMBER,"bold":True,"font":MONO}), ("  (828 vs 804, ×1.03)", {"color":TEXT})],
    [("PreVerify ≈ Verify", {"color":AMBER,"bold":True,"font":MONO}), ("  (197 vs 191, ×1.03)", {"color":TEXT})],
    [("Adapt ≈ Verify", {"color":TEXT}), (";  ", {"color":MUTED}), ("Ext is the cheapest op.", {"color":TEXT})],
], size=12.5, gap=8, marker="")
card(s, 8.05, 4.25, 4.55, 2.25, fill=CARD, line=LINE)
simple(s, 8.27, 4.38, 4.1, 0.4, "WHY", 11.5, MUTED, font=BODY, bold=True, spc=1.0)
bullets(s, 8.27, 4.78, 4.15, 1.6, [
    [("Folding the statement into the hash, ", {"color":TEXT}),
     ("H(pk, w+Y, M)", {"color":ACCENT,"font":MONO}),
     (", is free — one extra polynomial add.", {"color":TEXT})],
    [("Adding adaptor capability costs ", {"color":TEXT}),
     ("essentially nothing", {"color":ACCENT,"bold":True}),
     (" over the base scheme.", {"color":TEXT})],
], size=12.5, gap=9, marker="▸", marker_color=ACCENT)
notes(s, "The first signature-level metric. Eight operations timed. The visual deliberately pairs Sign with "
         "PreSign and Verify with PreVerify in amber: the bars are the same length. That is the result — the "
         "adaptor operations add about 3% over their base operations, because folding the statement into the "
         "Fiat-Shamir hash is a single polynomial addition. Adapt is about the cost of a Verify and Ext is "
         "the cheapest operation of all. This matches the paper's efficiency claim and sets up the contrast "
         "with the classical adaptor, whose overhead is 4x.")

# ----------------------------------------------------------------- 5 | ACCEPTANCE
s = base(prs, "Signature-level · rejection sampling", "Acceptance rate: measured, and it matches theory", n=5)
# big stat
statcard(s, 0.62, 2.0, 3.5, 1.7, "≈ 37%", "per-attempt acceptance\n(Sign 36.9% · PreSign 36.1%)", col=ACCENT, big_size=42)
statcard(s, 4.3, 2.0, 3.5, 1.7, "≈ 2.7", "attempts / signature\n(Sign 2.71 · PreSign 2.77)", col=ACCENT, big_size=42)
code_panel(s, 0.62, 3.95, 7.18, 1.55, title="MATCHES THE CLOSED FORM", accent=ACCENT, size=13.5, lines=[
    "(1 - k/g)^((n+l)*N) = (1 - 60/122880)^2048",
    "                   ~= e^(-1) ~= 36.8%   ->  ~2.72 attempts",
])
bullets(s, 0.62, 5.7, 7.2, 1.2, [
    [("Intrinsic to Fiat–Shamir-with-aborts; ", {"color":TEXT}),
     ("γ = κ·d·(n+ℓ)", {"color":ACCENT,"font":MONO}),
     (" governs MSIS hardness, not the accept rate.", {"color":TEXT})],
    [("Hint-free does ", {"color":TEXT}), ("not", {"color":AMBER,"bold":True}),
     (" worsen acceptance — optimised Dilithium rejects on more conditions, not fewer.", {"color":TEXT})],
], size=12.5, gap=8)
# methodology correction
card(s, 8.0, 2.0, 4.62, 4.5, fill=PANEL, line=AMBER)
simple(s, 8.22, 2.14, 4.2, 0.4, "A METHODOLOGY CORRECTION", 12, AMBER, font=BODY, bold=True, spc=0.8)
bullets(s, 8.22, 2.62, 4.25, 3.7, [
    [("An earlier bench ", {"color":TEXT}), ("estimated", {"color":CORAL,"bold":True}),
     (" retries from the timing ratio ", {"color":TEXT}),
     ("t_sign / t_verify", {"color":CORAL,"font":MONO}),
     (" → reported ~23% (~4.3 attempts).", {"color":TEXT})],
    [("That estimator is ", {"color":TEXT}), ("biased", {"color":CORAL,"bold":True}),
     (": a Sign attempt does n+ℓ=8 c·r products vs Verify's n=4 — so it over-counts.", {"color":TEXT})],
    [("The ", {"color":TEXT}), ("direct counter", {"color":ACCENT,"bold":True}),
     (" (~37%, ~2.7) supersedes it and agrees with the e⁻¹ line.", {"color":TEXT})],
    [("Prefer direct measurement to a proxy.", {"color":ACCENT,"italic":True,"bold":True})],
], size=12.5, gap=11, marker="▸", marker_color=AMBER)
notes(s, "The acceptance rate of the rejection-sampling loop. About 37% per attempt, about 2.7 attempts per "
         "signature, measured directly with the attempt counter. It matches the closed form: a signature is "
         "accepted only if all 2048 response coefficients land in range, giving (1-kappa/gamma)^2048, which "
         "is essentially e to the minus one. The amber panel is an honest methodology correction worth "
         "keeping: an earlier benchmark inferred the rate from a Sign/Verify timing ratio and got 23%, but "
         "that estimator is biased because a Sign attempt is dearer than a Verify, so it over-counts "
         "attempts. The direct counter replaces it — a small example of preferring measurement to a proxy.")

# ----------------------------------------------------------------- 6 | SIZES
s = base(prs, "Signature-level · sizes", "Object sizes: three distinct numbers — don't confuse them", n=6)
hdr = ["OBJECT", "IN-MEMORY sizeof", "PACKED (measured)", "PAPER (optimised)"]
colx = [0.0, 2.55, 5.0, 7.6]; colw = [2.55, 2.45, 2.6, 2.4]
tx, ty0 = 0.62, 2.0
rect(s, tx, ty0, 10.1, 0.5, fill=CARD2)
for j,h in enumerate(hdr):
    simple(s, tx+colx[j]+0.18, ty0, colw[j]-0.2, 0.5, h, 11.5,
           (ACCENT if j==2 else MUTED), font=BODY, bold=True, spc=0.6, anchor=MSO_ANCHOR.MIDDLE)
rows = [
    ("pk / statement Y", "4096 B", "2944 B", "—"),
    ("sk / witness y",   "8192 B", "512 B",  "—"),
    ("sig / pre-sig",    "9216 B", "4672 B", "~3210 B"),
]
rh = 0.66
for i,r in enumerate(rows):
    yy = ty0+0.5+i*rh
    rect(s, tx, yy, 10.1, rh, fill=(CARD if i%2==0 else PANEL))
    for j,val in enumerate(r):
        col = TEXT if j==0 else (ACCENT if j==2 else (MUTED if j!=1 else TEXT))
        simple(s, tx+colx[j]+0.18, yy, colw[j]-0.25, rh, val, 13.5, col,
               font=(BODY if j==0 else MONO), bold=(j==0 or j==2), anchor=MSO_ANCHOR.MIDDLE)
# packing breakdown
code_panel(s, 0.62, 4.7, 10.1, 1.75, title="WHERE THE PACKED SIG (4672 B) COMES FROM — bit-budget, not a formula guess", accent=ACCENT, size=12.5, lines=[
    [("challenge c : ternary, 2 bits/coeff  ->  256*2/8        =   64 B", {})],
    [("response  z : 18 bits/coeff (2*(g-k)+1 = 245641 < 2^18) ->  8*256*18/8 = 4608 B", {})],
    [("total       :                                              4672 B   (validated by test_serde)", {"color":ACCENT,"bold":True})],
])
# side note
card(s, 10.85, 2.0, 1.77, 2.48, fill=PANEL, line=LINE)
simple(s, 11.0, 2.12, 1.5, 1.0, "Compare LIKE FOR LIKE:", 10.5, AMBER, font=BODY, bold=True, line=1.0)
simple(s, 11.0, 2.95, 1.5, 1.5, "the paper's ~3210 B is a different (q≈2²⁴, hinted) scheme — use the PACKED column.",
       10.5, MUTED, font=BODY, line=1.05)
notes(s, "Sizes, and the single most common confusion. There are three different numbers and they must not "
         "be mixed. In-memory sizeof stores a full int32 per coefficient — that is the 9216-byte figure and "
         "it is an artefact of the data structure. The packed column is the real serialised wire size, 4672 "
         "bytes, validated byte-for-byte by the serde test. The paper's ~3210 bytes is a different, "
         "optimised scheme at a different modulus with a hint vector, so it is not directly comparable; the "
         "honest comparison for our scheme is always the packed column. The bottom panel shows the packed "
         "signature is a real bit-budget, not a guess.")

# ----------------------------------------------------------------- 7 | 2x2 MATRIX
s = base(prs, "The 2×2 matrix", "Four schemes, one machine: basic/exotic × classical/PQ", n=7)
simple(s, 0.62, 1.72, 12.0, 0.35, "µs/op, 2000 iters, same machine & compiler (bench_classical + bench_las3)", 12, MUTED, font=BODY, italic=True)
cols = [
    ("OPERATION", MUTED, 0.62, 2.55),
    ("ECDSA", DIM, 3.2, 1.85),
    ("ECDSA-adaptor", CORAL, 5.1, 2.35),
    ("Dilithium-3", VIOLET, 7.5, 2.05),
    ("LAS", ACCENT, 9.6, 3.0),
]
ty0 = 2.18
rect(s, 0.62, ty0, 11.98, 0.52, fill=CARD2)
for (lab,col,cx,cw) in cols:
    simple(s, cx+0.15, ty0, cw-0.2, 0.52, lab, 12.5, col, font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
data = [
    ("KeyGen / stmt", "31", "31", "162", "78"),
    ("Sign", "41", "—", "642", "804"),
    ("Verify", "62", "—", "155", "191"),
    ("PreSign", "—", "189", "—", "828"),
    ("PreVerify", "—", "244", "—", "197"),
    ("Adapt", "—", "3", "—", "203"),
    ("Ext", "—", "35", "—", "68"),
]
rh = 0.50
for i,row in enumerate(data):
    yy = ty0+0.52+i*rh
    rect(s, 0.62, yy, 11.98, rh, fill=(CARD if i%2==0 else PANEL))
    # highlight LAS column
    rect(s, 9.6, yy, 3.0, rh, fill=(CARD2 if i%2==0 else CARD))
    for k,(lab,col,cx,cw) in enumerate(cols):
        val = row[k]
        c = TEXT if k==0 else (DIM if val=="—" else (ACCENT if k==4 else TEXT))
        simple(s, cx+0.15, yy, cw-0.2, rh, val, 13, c,
               font=(BODY if k==0 else MONO), bold=(k==0 or k==4), anchor=MSO_ANCHOR.MIDDLE)
simple(s, 0.62, 6.5, 12.0, 0.4,
       "Everything stays sub-millisecond on commodity hardware — the PQ column is feasible, not exotic-slow.",
       13, ACCENT, font=BODY, bold=True, italic=True)
notes(s, "The centrepiece table: the project's own 2x2 framing turned into data. Four schemes on one "
         "machine — classical-basic ECDSA, classical-exotic ECDSA-adaptor, PQ-basic Dilithium-3, and "
         "PQ-exotic LAS. The LAS column is highlighted. Read it two ways, which are the next two slides: "
         "down the columns for the cost of going post-quantum, and within each adaptor scheme for the cost "
         "of going exotic. The headline from the raw table is simply that every operation, in every scheme, "
         "is sub-millisecond — post-quantum here is feasible, not prohibitively slow.")

# ----------------------------------------------------------------- 8 | FINDING 1
s = base(prs, "Finding 1", "The price of post-quantum is communication, not computation", n=8, accent=AMBER)
# left: size ratios (tall bars), right: time ratios (short bars), same axis
base_y = 5.7; max_h = 3.0; maxr = 90
simple(s, 0.9, 1.95, 5.6, 0.35, "COMMUNICATION — size ratio (LAS ÷ ECDSA)", 12, CORAL, font=BODY, bold=True, spc=0.4)
sz = [("pk",89,CORAL),("sig",73,CORAL),("pre-sig",29,CORAL),("sk",16,CORAL)]
for i,(lab,v,col) in enumerate(sz):
    cx = 1.0 + i*1.35
    vbar(s, cx, base_y, 0.85, v, maxr, max_h, col, toplabel=f"×{v}", botlabel=lab, label_w=1.2)
simple(s, 7.4, 1.95, 5.2, 0.35, "COMPUTATION — time ratio (LAS ÷ ECDSA)", 12, ACCENT, font=BODY, bold=True, spc=0.4)
tm = [("Sign",19.5,ACCENT),("PreSign",4.4,ACCENT),("Verify",3.1,ACCENT)]
for i,(lab,v,col) in enumerate(tm):
    cx = 7.55 + i*1.55
    vbar(s, cx, base_y, 0.95, v, maxr, max_h, col, toplabel=f"×{v}", botlabel=lab, label_w=1.4)
# divider + baseline
rect(s, 6.95, 2.0, 0.014, 3.95, fill=LINE)
rect(s, 0.9, base_y, 11.7, 0.018, fill=LINE)
# takeaway
card(s, 0.62, 6.15, 12.0, 0.95, fill=CARD2, line=AMBER)
bullets(s, 0.85, 6.24, 11.6, 0.8, [
    [("Sizes grow ×16–×89; times grow ×3–×20 and stay sub-ms. ", {"color":TEXT,"bold":True}),
     ("For a blockchain the size column is the binding constraint — on-chain the 4672-B sig's calldata alone (74,476 gas) exceeds the entire classical claim.", {"color":TEXT})],
], size=13, gap=2, marker="")
notes(s, "Finding one, read down the columns of the matrix. The bars use one shared vertical scale. On the "
         "left, the size ratios of LAS over ECDSA: public key 89 times, signature 73 times, pre-signature "
         "29, secret key 16. On the right, the time ratios on the same scale: Sign 19.5, PreSign 4.4, Verify "
         "3.1 — visibly tiny by comparison, and everything stays sub-millisecond. So the cost of going "
         "post-quantum is overwhelmingly communication, not computation. That matters enormously for "
         "blockchain, where bytes are the binding constraint — as the gas slides will show, the signature's "
         "calldata alone outweighs the whole classical settlement.")

# ----------------------------------------------------------------- 9 | FINDING 2
s = base(prs, "Finding 2 · LAS's headline win", "The adaptor overhead is inverted", n=9)
simple(s, 0.62, 1.78, 12.0, 0.4, "Cost of the adaptor operation relative to its own base operation (× base):",
       14, TEXT, font=BODY, bold=True)
bars = [
    ("Classical PreSign / Sign",   4.6, CORAL),
    ("LAS PreSign / Sign",         1.03, ACCENT),
    ("Classical PreVerify / Verify", 3.9, CORAL),
    ("LAS PreVerify / Verify",      1.03, ACCENT),
]
by0, bh, bgap, track, maxv = 2.45, 0.42, 0.28, 6.3, 5.0
for i,(lab,v,col) in enumerate(bars):
    yy = by0 + i*(bh+bgap)
    simple(s, 0.62, yy-0.02, 4.0, bh+0.04, lab, 12.5, (col if col!=CORAL else TEXT), font=BODY, bold=(col==ACCENT), anchor=MSO_ANCHOR.MIDDLE)
    hbar(s, 4.75, yy, track, bh, v, maxv, col, valstr=f"×{v}", val_color=col)
# 1x reference line
ref_x = 4.75 + track*(1.0/maxv)
rect(s, ref_x, by0-0.12, 0.014, (bh+bgap)*4-0.16, fill=MUTED)
simple(s, ref_x-0.45, by0-0.42, 1.0, 0.3, "×1 (free)", 10, MUTED, font=BODY, italic=True, align=CENTER)
# callouts
card(s, 0.62, 5.5, 5.95, 1.45, fill=CARD2, line=ACCENT)
bullets(s, 0.85, 5.62, 5.6, 1.2, [
    [("Classical pays for a ", {"color":TEXT}), ("DLEQ proof", {"color":CORAL,"bold":True}),
     (" in the pre-signature; LAS folds Y into the hash for free.", {"color":TEXT})],
], size=13, gap=2, marker="")
card(s, 6.7, 5.5, 5.92, 1.45, fill=CARD2, line=ACCENT)
bullets(s, 6.93, 5.62, 5.55, 1.2, [
    [("LAS PreVerify (197 µs) is ", {"color":TEXT}),
     ("absolutely faster", {"color":ACCENT,"bold":True}),
     (" than classical ECDSA-adaptor PreVerify (244 µs) — same machine.", {"color":TEXT})],
], size=13, gap=2, marker="")
notes(s, "Finding two, read within each adaptor scheme — and it is LAS's headline win. In the classical "
         "ECDSA-adaptor, the adaptor functionality is expensive relative to its own base: PreSign is 4.6 "
         "times Sign and PreVerify is 3.9 times Verify, because the pre-signature must carry and verify a "
         "DLEQ proof. In LAS, PreSign and PreVerify are both about 1.03 times their base — essentially free "
         "— because the statement just folds into the Fiat-Shamir hash. The overhead structure is inverted. "
         "And strikingly, LAS PreVerify in absolute terms is faster than the classical adaptor's PreVerify "
         "on the same machine. The exotic PQ scheme beats the exotic classical scheme on adaptor overhead.")

# ----------------------------------------------------------------- 10 | APP: SWAP PAYLOAD
s = base(prs, "Application-level · atomic swap", "Swap payload: what actually crosses the wire", n=10, accent=AMBER)
simple(s, 0.62, 1.75, 12.0, 0.35, "2 parties, 2 chains, scriptless (bench_app3) — actual serialised bytes",
       12, MUTED, font=BODY, italic=True)
# flow of three big numbers
phases = [
    ("OFF-CHAIN", "12,288 B", "Y + σ̂_A + σ̂_B\n(2944 + 4672 + 4672)", AMBER),
    ("SETTLEMENT", "9,344 B", "σ_A, σ_B published\n(2 × 4672)", ACCENT),
    ("+ ESCROWED Y", "15,232 B", "settlement + 2 × Y", MUTED),
]
pw, ph = 3.85, 2.2; px0, py = 0.62, 2.35; gap = 0.28
for i,(t,big,sub,col) in enumerate(phases):
    cx = px0 + i*(pw+gap)
    card(s, cx, py, pw, ph, fill=CARD, line=col)
    rect(s, cx, py, pw, 0.5, fill=PANEL); rect(s, cx, py+0.30, pw, 0.2, fill=PANEL)
    simple(s, cx+0.22, py, pw-0.4, 0.5, t, 12, col, font=BODY, bold=True, spc=0.8, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, cx+0.22, py+0.62, pw-0.4, 0.8, big, 34, col, font=HEAD, bold=True)
    tb,tf = textbox(s, cx+0.22, py+1.45, pw-0.4, 0.7)
    for k,row in enumerate(sub.split("\n")):
        p = para(tf, first=(k==0), space_after=1, line=1.0); _run(p, row, 12, TEXT, font=BODY)
    if i<2:
        simple(s, cx+pw+0.01, py, gap, ph, "→", 22, DIM, font=BODY, bold=True, align=CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 0.62, 4.85, 12.0, 1.8, [
    [("Only the ", {"color":TEXT}), ("two adapted signatures", {"color":ACCENT,"bold":True}),
     (" would actually be published on a real chain — each is a single ordinary-looking LAS signature.", {"color":TEXT})],
    [("End-to-end signing work (2× PreSign + 2× Adapt + Ext) is ", {"color":TEXT}),
     ("a few milliseconds", {"color":ACCENT,"bold":True}),
     (", dominated by the two rejection-sampled pre-signs.", {"color":TEXT})],
    [("Harness re-asserts the fairness invariant: ", {"color":TEXT}),
     ("adapted sigs verify, pre-sigs do not.", {"color":AMBER,"bold":True})],
], size=13.5, gap=10)
notes(s, "Moving to the application level — the second benchmark type. For the atomic cross-chain swap, the "
         "off-chain phase exchanges three messages totalling about 12 kilobytes: the statement plus two "
         "pre-signatures. The on-chain settlement footprint is the two adapted signatures, about 9.3 "
         "kilobytes, and only those would actually hit a real chain. Each adapted signature is a single, "
         "ordinary-looking LAS signature. End-to-end signing is a few milliseconds, dominated by the two "
         "rejection-sampled pre-signs. The harness asserts the fairness invariant every run: the adapted "
         "signatures verify and the pre-signatures do not.")

# ----------------------------------------------------------------- 11 | APP: AMHL vs K
s = base(prs, "Application-level · multi-hop (AMHL)", "Cost as a function of path length K", n=11, accent=AMBER)
simple(s, 0.62, 1.72, 12.0, 0.35, "40 routes per K, mode 3 (bench_app3) — settlement bytes grow linearly; signing stays flat",
       12, MUTED, font=BODY, italic=True)
# vertical bars: settlement bytes for K=1,2,4,6,8
Ks = [(1,4672,2.60,1),(2,9344,2.30,2),(4,18688,2.73,4),(6,28032,2.95,6),(8,37376,2.91,7)]
base_y = 5.35; max_h = 2.7; maxv = 38000
for i,(K,by,att,sn) in enumerate(Ks):
    cx = 1.1 + i*1.18
    vbar(s, cx, base_y, 0.8, by, maxv, max_h, ACCENT,
         toplabel=f"{by//1000}k", botlabel=f"K={K}", tl_color=ACCENT, label_w=1.1)
rect(s, 0.95, base_y, 6.4, 0.018, fill=LINE)
simple(s, 0.95, 2.0, 6.4, 0.3, "Settlement footprint (B)  =  K × 4672", 12, ACCENT, font=BODY, bold=True)
# right: the three findings
findings = [
    ("Linear in K", "K adapted sigs + K statements — no super-linear blow-up.", ACCENT),
    ("Witness norm grows", "‖s_j‖∞ ≤ j ≤ K (sum of j ternary vectors) — the reason every hop pre-signs at γ−κ−K.", AMBER),
    ("γ−κ−K is free", "K=1→8 shrinks the accept band by only 7/(γ−κ) ≈ 0.0057% — attempts/presig stays ≈2.7–3.0.", ACCENT),
]
for i,(t,d,col) in enumerate(findings):
    yy = 2.15 + i*1.45
    card(s, 7.65, yy, 4.97, 1.3, fill=CARD, line=LINE)
    rect(s, 7.65, yy, 0.09, 1.3, fill=col)
    simple(s, 7.9, yy+0.14, 4.6, 0.4, t, 14, col, font=BODY, bold=True)
    simple(s, 7.9, yy+0.55, 4.6, 0.7, d, 12, TEXT, font=BODY, line=1.06)
notes(s, "The multi-hop AMHL payment, costed as a function of path length K. Three findings. First, the "
         "settlement footprint is exactly linear in K — K adapted signatures plus K statements, no "
         "super-linear blow-up. Second, the witness norm grows with the hop index, bounded by j and "
         "therefore by K, because each cumulative witness is a sum of j ternary vectors; that is the "
         "knowledge gap made concrete and the precise reason every hop must pre-sign at the tighter "
         "gamma-kappa-K bound. Third, and slightly counter-intuitively, that tightening is free: going from "
         "one hop to eight shrinks the acceptance band by about half a hundredth of a percent, so the "
         "attempts per pre-sign stay flat. AMHL adds no per-hop signing penalty beyond the unavoidable K "
         "pre-signatures.")

# ----------------------------------------------------------------- 12 | ON-CHAIN GAS SWAP
s = base(prs, "On-chain · real Solidity swap", "EVM gas: claim cost, classical vs LAS", n=12, accent=CORAL)
simple(s, 0.62, 1.72, 12.0, 0.35, "Foundry local EVM, shared HTLC escrow — gas is deterministic (evm/, forge --gas-report)",
       12, MUTED, font=BODY, italic=True)
# claim comparison bars
simple(s, 0.62, 2.15, 7.0, 0.35, "CLAIM (settlement + signature check) — gas", 12, MUTED, font=BODY, bold=True, spc=0.6)
gb = [
    ("Classical claim  (full ecrecover verify)", 75709, ACCENT, "75,709"),
    ("LAS claim FLOOR  (no real verification)", 208400, CORAL, "208,400"),
]
by0, bh, track, maxv = 2.65, 0.5, 4.4, 215000
for i,(lab,v,col,vs) in enumerate(gb):
    yy = by0 + i*0.82
    simple(s, 0.62, yy-0.24, 6.5, 0.3, lab, 11.5, TEXT, font=BODY, bold=True)
    hbar(s, 0.66, yy, track, bh, v, maxv, col, valstr=vs, val_color=col)
# calldata callout
card(s, 0.62, 4.45, 6.0, 2.05, fill=CARD2, line=CORAL)
simple(s, 0.84, 4.57, 5.6, 0.4, "97% OF IT IS JUST THE BYTES", 12, CORAL, font=BODY, bold=True, spc=0.6)
bullets(s, 0.84, 5.0, 5.6, 1.45, [
    [("The 4672-B sig = 4649 non-zero + 23 zero bytes → ", {"color":TEXT}),
     ("74,476 gas of calldata alone", {"color":CORAL,"bold":True}),
     (" (16/non-zero, 4/zero).", {"color":TEXT})],
    [("That ", {"color":TEXT}), ("exceeds the entire classical claim", {"color":CORAL,"bold":True}),
     (" (75,709) — which includes real verification.", {"color":TEXT})],
], size=12.5, gap=9, marker="▸", marker_color=CORAL)
# right: full table
simple(s, 7.0, 2.15, 5.6, 0.35, "FULL GAS REPORT", 12, MUTED, font=BODY, bold=True, spc=0.8)
gt = [("","Classical","LAS"),("fund","180,285","139,568"),
      ("claim","75,709","208,400"),("refund","39,330","39,330"),("deploy","715,257","—")]
ty0=2.6; rh=0.6
for i,(a,b,c) in enumerate(gt):
    yy=ty0+i*rh
    if i==0:
        rect(s,7.0,yy,5.62,rh,fill=CARD2)
        simple(s,7.2,yy,2.0,rh,a,12,MUTED,font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
        simple(s,9.0,yy,1.8,rh,b,12,ACCENT,font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=CENTER)
        simple(s,10.8,yy,1.7,rh,c,12,CORAL,font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=CENTER)
    else:
        rect(s,7.0,yy,5.62,rh,fill=(CARD if i%2==1 else PANEL))
        simple(s,7.2,yy,2.0,rh,a,12.5,TEXT,font=BODY,bold=True,anchor=MSO_ANCHOR.MIDDLE)
        simple(s,9.0,yy,1.8,rh,b,12.5,TEXT,font=MONO,anchor=MSO_ANCHOR.MIDDLE,align=CENTER)
        simple(s,10.8,yy,1.7,rh,c,12.5,(CORAL if a=="claim" else TEXT),font=MONO,bold=(a=="claim"),anchor=MSO_ANCHOR.MIDDLE,align=CENTER)
simple(s, 7.0, 5.7, 5.6, 0.7, "LAS settlement floor = 2.75× the full classical claim — and the floor does no verification at all.",
       12.5, AMBER, font=BODY, bold=True, italic=True, line=1.1)
notes(s, "On-chain, on a real EVM. The escrow contract is signature-scheme-agnostic; only the claim-time "
         "verification differs, so the gas report isolates the price of post-quantum on-chain. The classical "
         "claim verifies the adapted ECDSA signature natively with ecrecover for 75,709 gas. The LAS claim "
         "can only charge the unavoidable floor — calldata for the real 4672-byte signature plus one keccak "
         "pass — and that floor is already 208,400 gas, 2.75 times the full classical claim, while doing no "
         "verification. The reason is again communication: the signature's calldata alone, 74,476 gas, is "
         "97% of the marginal cost and on its own exceeds the entire classical claim.")

# ----------------------------------------------------------------- 13 | NATIVE VERIFY EXPERIMENT
s = base(prs, "On-chain · the cost experiment", "What native LAS verification would actually cost", n=13, accent=CORAL)
simple(s, 0.62, 1.72, 12.0, 0.35, "LASVerifyCost.t.sol — a gas-faithful probe of one las_verify (op-budget identical to ref/las.c)",
       12, MUTED, font=BODY, italic=True)
# breakdown table
bd = [
    ("forward NTT × 12", "4,537,776", "measured"),
    ("inverse NTT × 8", "3,374,048", "measured"),
    ("pointwise × 20", "940,500", "measured"),
    ("coeff passes × ~40", "≈1,227,720", "measured"),
    ("arithmetic (verifyArith)", "10,080,044", "MEASURED"),
    ("SHAKE256 challenge × 64 perm.", "1,920,000", "calculated"),
    ("one native las_verify", "≈12,000,044", "TOTAL"),
]
tx, ty0, rh = 0.62, 2.15, 0.535
for i,(a,b,c) in enumerate(bd):
    yy = ty0 + i*rh
    emph = a.startswith("arithmetic") or a.startswith("one native")
    rect(s, tx, yy, 7.05, rh-0.06, fill=(CARD2 if emph else (CARD if i%2==0 else PANEL)),
         line=(ACCENT if a.startswith("arithmetic") else (CORAL if a.startswith("one") else None)))
    simple(s, tx+0.2, yy, 3.6, rh-0.06, a, 12.5, (CORAL if a.startswith("one") else (ACCENT if emph else TEXT)),
           font=BODY, bold=emph, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, tx+3.85, yy, 1.95, rh-0.06, b, 13, (CORAL if a.startswith("one") else TEXT),
           font=MONO, bold=emph, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    simple(s, tx+5.95, yy, 1.0, rh-0.06, c, 9.5, (AMBER if c in ("MEASURED","TOTAL") else MUTED),
           font=BODY, bold=(c in ("MEASURED","TOTAL")), anchor=MSO_ANCHOR.MIDDLE, align=CENTER)
# right: the verdict
statcard(s, 8.0, 2.15, 4.62, 1.5, "≈ 12 M gas", "one native verification\n(measured + calculated)", col=CORAL, big_size=34)
card(s, 8.0, 3.8, 4.62, 1.35, fill=CARD, line=LINE)
bullets(s, 8.22, 3.92, 4.2, 1.1, [
    [("≈ 158×", {"color":CORAL,"bold":True,"font":MONO}), (" the classical ecrecover claim", {"color":TEXT})],
    [("≈ 40%", {"color":CORAL,"bold":True,"font":MONO}), (" of a single 30 M block", {"color":TEXT})],
], size=13, gap=8, marker="▸", marker_color=CORAL)
card(s, 8.0, 5.25, 4.62, 1.25, fill=PANEL, line=AMBER)
bullets(s, 8.22, 5.36, 4.2, 1.05, [
    [("Retraction:", {"color":AMBER,"bold":True}),
     (' the earlier "exceeds the block limit" claim was wrong — with native mulmod it ', {"color":TEXT}),
     ("fits in a block", {"color":TEXT,"bold":True}),
     (". The real barrier is cost + missing precompiles.", {"color":TEXT})],
], size=12, gap=2, marker="")
notes(s, "Because I claimed on-chain verification was infeasible, I measured it rather than asserting it. "
         "The probe runs the exact op-budget of one las_verify on the EVM — 12 forward NTTs, 8 inverse NTTs, "
         "20 pointwise products — and prices it; EVM opcode gas is independent of operand values, so a "
         "structurally identical kernel has the same gas as a correct verifier. The arithmetic measures at "
         "10.08 million gas, and the first four rows independently reconcile to that total. Adding the "
         "SHAKE256 challenge, about 64 Keccak permutations, gives roughly 12 million gas total — about 158 "
         "times the classical claim and 40% of a 30-million-gas block. Two honest consequences: an earlier "
         "draft's claim that it 'exceeds the block limit' was wrong and is retracted, because native mulmod "
         "makes the arithmetic fit; the real barriers are economics and missing precompiles. That is "
         "precisely the case for a PQ precompile or a zk proof — the poqeth conclusion, here quantified for "
         "the exotic case.")

# ----------------------------------------------------------------- 14 | CORRECTNESS / ROBUSTNESS
s = base(prs, "Correctness & robustness", "The metrics that have no margin of error", n=14, accent=VIOLET)
cards = [
    ("100%", "correctness over 1000 iterations × modes 2/3/5 — PreVerify/Verify/Adapt/Ext all hold every run", ACCENT),
    ("4672 / 4672", "single-byte signature flips rejected — exhaustive tamper test via the packed verifier", ACCENT),
    ("1 pinned digest", "SHAKE256 KAT byte-identical across runs & machines:  f7fc40f0b7752caf…", ACCENT),
]
for i,(big,desc,col) in enumerate(cards):
    yy = 2.05 + i*1.5
    card(s, 0.62, yy, 6.4, 1.32, fill=CARD, line=col)
    rect(s, 0.62, yy, 0.10, 1.32, fill=col)
    simple(s, 0.9, yy+0.16, 6.0, 0.6, big, 28, col, font=HEAD, bold=True)
    simple(s, 0.9, yy+0.78, 6.05, 0.5, desc, 12.5, TEXT, font=BODY, line=1.05)
# the tripwire highlight
card(s, 7.25, 2.05, 5.37, 4.32, fill=PANEL, line=AMBER)
simple(s, 7.47, 2.2, 5.0, 0.4, "THE TRIPWIRE (statement binding)", 12.5, AMBER, font=BODY, bold=True, spc=0.6)
bullets(s, 7.47, 2.7, 5.05, 3.5, [
    [("Feed a pre-signature σ̂ to ordinary Verify → it ", {"color":TEXT}),
     ("must fail.", {"color":CORAL,"bold":True})],
    [("Why: Verify recomputes ", {"color":TEXT}),
     ("H(pk, w, M)", {"color":ACCENT,"font":MONO}),
     (", but the pre-sig's challenge is ", {"color":TEXT}),
     ("H(pk, w+Y, M)", {"color":AMBER,"font":MONO}),
     (" — different with overwhelming probability.", {"color":TEXT})],
    [("This proves the ", {"color":TEXT}), ("+Y is genuinely in the hash", {"color":ACCENT,"bold":True}),
     (" and the pre-sig is unspendable — asserted every iteration.", {"color":TEXT})],
    [("Same property at the multi-hop level: adapting hop 1 with s_K ⇒ Verify fails ", {"color":TEXT}),
     ("(wormhole resistance).", {"color":AMBER,"bold":True})],
], size=13, gap=13, marker="▸", marker_color=AMBER)
notes(s, "Correctness is not a benchmark you optimise — it is pass/fail, and these have no margin. A "
         "thousand iterations across all three Dilithium modes, 100% correct on every adaptor property. An "
         "exhaustive tamper test: every one of the 4672 single-byte flips of a packed signature is rejected. "
         "And a pinned known-answer-test digest that is byte-identical across runs and machines. The right "
         "panel is the most important assertion in the whole suite: the tripwire. Feeding a pre-signature to "
         "ordinary Verify must fail, because Verify hashes w while the pre-sig committed to w+Y. That single "
         "failing check is the executable proof that the statement really is bound into the signature and "
         "that pre-signatures are unspendable — and at the multi-hop level it becomes wormhole resistance.")

# ----------------------------------------------------------------- 15 | SYNTHESIS
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG)
rect(s, 0.0, 0.0, 0.18, SH, fill=ACCENT)
dot_grid(s, 11.55, 0.45, 4, 4, step=0.34, r=0.05, color=LATTICE)
simple(s, 0.7, 0.62, 11.0, 0.35, "SYNTHESIS", 13, ACCENT, font=BODY, bold=True, spc=2.6)
simple(s, 0.68, 1.04, 11.8, 0.9, "Three findings the data delivers", 33, TEXT, font=HEAD, bold=True)
finds = [
    ("Price of PQ = communication, not computation",
     "sizes ×16–89 and on-chain calldata (74k gas) dominate; times stay ×3–20 and sub-ms.", CORAL),
    ("LAS's adaptor overhead is ≈0 — and inverted vs classical",
     "PreSign ≈ Sign, PreVerify ≈ Verify (×1.03) vs the classical ~4× DLEQ; LAS PreVerify even absolutely faster.", ACCENT),
    ("The swap runs end-to-end; native on-chain verify is the open piece",
     "≈12 M gas (158× classical, 40% of a block) — prohibitive but NOT over the block limit; wants a precompile or zk.", AMBER),
]
for i,(t,d,col) in enumerate(finds):
    yy = 2.2 + i*1.18
    card(s, 0.7, yy, 11.9, 1.05, fill=CARD, line=None)
    rect(s, 0.7, yy, 0.10, 1.05, fill=col)
    o=rect(s, 0.95, yy+1.05/2-0.30, 0.6, 0.6, fill=PANEL, line=col, lw=1.5, shape=MSO_SHAPE.OVAL)
    tf=o.text_frame;_set_margins(tf,0.0);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    pp=tf.paragraphs[0];pp.alignment=CENTER;_run(pp,str(i+1),20,col,font=HEAD,bold=True)
    simple(s, 1.75, yy+0.13, 10.6, 0.45, t, 16.5, TEXT, font=BODY, bold=True)
    simple(s, 1.75, yy+0.58, 10.6, 0.45, d, 12.5, MUTED, font=BODY, line=1.05)
simple(s, 0.7, 5.95, 11.9, 0.85,
       "Aligned with the goals: two benchmark types × two baselines, the 2×2 quadrant quantified, every claim measured — and every caveat named.",
       14, ACCENT, font=BODY, bold=True, italic=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 15)
notes(s, "The synthesis. Three findings the evaluation delivers. One: the price of post-quantum is "
         "communication, not computation — sizes and calldata dominate while times stay small and sub-ms. "
         "Two: LAS's adaptor overhead is essentially zero and structurally inverted versus the classical "
         "adaptor's roughly fourfold DLEQ cost, with LAS PreVerify even absolutely faster. Three: the swap "
         "protocol runs end-to-end, and the one genuinely open piece is native on-chain verification at "
         "about 12 million gas — prohibitively expensive but, measured, not over the block limit, which "
         "points to a precompile or a zk proof. All of this satisfies the supervisor's brief: two benchmark "
         "types, two baselines, the 2x2 framing turned into data, every claim measured and every caveat "
         "named.")

# ----------------------------------------------------------------- write
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LAS_Evaluation_Analysis.pptx")
prs.save(out)
print("SAVED:", out)
print("slides:", len(prs.slides._sldIdLst))
