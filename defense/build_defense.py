# -*- coding: utf-8 -*-
"""
Build the LAS-on-Dilithium DESIGN DEFENCE deck (.pptx).

A viva-grade slide deck that defends *why* LAS (eprint 2020/845, Algorithm 2)
was implemented on top of the CRYSTALS-Dilithium reference repo the way it was.
Each slide answers one design-decision question an examiner would ask, grounded
in the paper's algorithm and the actual C implementation (ref/las.{c,h}).

Reproducible artefact: re-run to regenerate the .pptx.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------- palette
BG      = RGBColor(0x0D, 0x10, 0x30)   # deep indigo-navy (dominant)
PANEL   = RGBColor(0x16, 0x1B, 0x45)   # panel
CARD    = RGBColor(0x1E, 0x24, 0x57)   # card
CARD2   = RGBColor(0x26, 0x2D, 0x66)   # lighter card
ACCENT  = RGBColor(0x49, 0xE3, 0xCE)   # cyan/mint  -> reuse / correct / positive
AMBER   = RGBColor(0xFF, 0xB4, 0x54)   # amber      -> a defended decision / highlight
CORAL   = RGBColor(0xFF, 0x6B, 0x6B)   # coral      -> failure mode / "wrong"
TEXT    = RGBColor(0xEC, 0xEE, 0xFF)   # near-white body
MUTED   = RGBColor(0x9A, 0xA0, 0xD0)   # muted lavender
DIM     = RGBColor(0x70, 0x78, 0xB0)   # dimmer
LINE    = RGBColor(0x2E, 0x35, 0x6F)   # subtle divider
CODEBG  = RGBColor(0x0A, 0x0D, 0x28)   # code panel bg
LATTICE = RGBColor(0x29, 0x31, 0x68)   # faint lattice dots

HEAD = "Georgia"
BODY = "Calibri"
MONO = "Consolas"

EMU_IN = 914400
SW, SH = 13.333, 7.5

# ----------------------------------------------------------------------------- helpers
def set_bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def no_shadow(sh):
    sh.shadow.inherit = False

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    no_shadow(sh)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sh.adjustments[0] = radius
        except Exception: pass
    return sh

def _set_margins(tf, m=0.06):
    tf.margin_left = Inches(m); tf.margin_right = Inches(m)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)

def _run(p, text, size, color, font=BODY, bold=False, italic=False, spc=None):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.name = font; f.bold = bold; f.italic = italic
    f.color.rgb = color
    if spc is not None:
        rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(int(spc * 100)))
    return r

def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True, m=0.06):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    _set_margins(tf, m)
    tf.paragraphs[0].alignment = align
    return tb, tf

def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    return p

def simple(s, x, y, w, h, text, size, color, font=BODY, bold=False, italic=False,
           align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, line=None, m=0.06):
    tb, tf = textbox(s, x, y, w, h, anchor=anchor, align=align, m=m)
    p = para(tf, first=True, align=align, space_after=0, line=line)
    _run(p, text, size, color, font=font, bold=bold, italic=italic, spc=spc)
    return tb

def bullets(s, x, y, w, h, items, size=14.5, gap=7, color=TEXT, marker_color=ACCENT,
            marker="▸", font=BODY, line=1.0, anchor=MSO_ANCHOR.TOP):
    """items: each item is a string OR a list of run-tuples (text, opts-dict)."""
    tb, tf = textbox(s, x, y, w, h, anchor=anchor)
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0), space_after=gap, line=line)
        if marker:
            _run(p, marker + "  ", size, marker_color, font=font, bold=True)
        if isinstance(it, str):
            _run(p, it, size, color, font=font)
        else:
            for (t, o) in it:
                _run(p, t, o.get("size", size), o.get("color", color),
                     font=o.get("font", font), bold=o.get("bold", False),
                     italic=o.get("italic", False))
    return tb

def dot_grid(s, x0, y0, cols, rows, step=0.30, r=0.040, color=LATTICE, link=False):
    for i in range(cols):
        for j in range(rows):
            o = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(x0 + i * step), Inches(y0 + j * step),
                                   Inches(r), Inches(r))
            no_shadow(o); o.line.fill.background()
            o.fill.solid(); o.fill.fore_color.rgb = color

def kicker_title(s, kicker, title, accent=ACCENT, title_size=30, ty=0.74, tw=12.1):
    simple(s, 0.62, 0.42, 11.0, 0.32, kicker.upper(), 12.5, accent, font=BODY, bold=True, spc=2.4)
    simple(s, 0.6, ty, tw, 0.95, title, title_size, TEXT, font=HEAD, bold=True)

def footer(s, n):
    simple(s, 0.62, 7.06, 9.0, 0.3, "LAS on CRYSTALS-Dilithium  ·  Design Defence  ·  eprint 2020/845",
           9, DIM, font=BODY)
    simple(s, 12.0, 7.06, 0.9, 0.3, f"{n:02d}", 9.5, MUTED, font=MONO, bold=True, align=PP_ALIGN.RIGHT)

def base(prs, kicker, title, accent=ACCENT, n=0, motif=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    if motif:
        dot_grid(s, 12.05, 0.30, 4, 3)          # subtle corner lattice
    kicker_title(s, kicker, title, accent=accent)
    footer(s, n)
    return s

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def code_panel(s, x, y, w, h, lines, size=12.5, title=None, accent=ACCENT,
               fill=CODEBG, border=LINE):
    rect(s, x, y, w, h, fill=fill, line=border, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    pad = 0.16
    yy = y + pad
    if title:
        simple(s, x + pad, yy, w - 2 * pad, 0.3, title, 11, accent, font=BODY, bold=True, spc=1.5)
        yy += 0.34
    tb, tf = textbox(s, x + pad, yy, w - 2 * pad, h - (yy - y) - pad)
    for i, ln in enumerate(lines):
        p = para(tf, first=(i == 0), space_after=3, line=1.04)
        if isinstance(ln, str):
            _run(p, ln if ln else " ", size, TEXT, font=MONO)
        else:
            for (t, o) in ln:
                _run(p, t, o.get("size", size), o.get("color", TEXT), font=o.get("font", MONO),
                     bold=o.get("bold", False), italic=o.get("italic", False))
    return tb

def chip(s, x, y, w, text, color, fill, h=0.34, size=11):
    sh = rect(s, x, y, w, h, fill=fill, line=color, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sh.text_frame; _set_margins(tf, 0.05); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, text, size, color, font=BODY, bold=True)
    return sh

def card(s, x, y, w, h, fill=CARD, line=None, radius=0.045):
    return rect(s, x, y, w, h, fill=fill, line=line, lw=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)

# ============================================================================= build
prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
prs.author = "Royce Steven"
prs.title = "LAS on Dilithium - Design Defence"

# ----------------------------------------------------------------- 1 | TITLE
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
rect(s, 0, 0, SW, SH, fill=BG)
# decorative lattice field (top-right)
dot_grid(s, 9.2, 0.5, 11, 6, step=0.36, r=0.05, color=LATTICE)
# accent band
rect(s, 0.0, 0.0, 0.18, SH, fill=ACCENT)
simple(s, 0.7, 1.35, 11.5, 0.4, "MSc CYBERSECURITY THESIS · UNIVERSITY OF MANCHESTER",
       13, ACCENT, font=BODY, bold=True, spc=2.6)
simple(s, 0.68, 1.95, 11.7, 1.7,
       "Defending the Design", 52, TEXT, font=HEAD, bold=True)
simple(s, 0.7, 2.95, 11.8, 1.2,
       "Why LAS is built on CRYSTALS-Dilithium the way it is", 26, MUTED, font=HEAD, italic=True)
# rule
rect(s, 0.72, 4.25, 6.6, 0.022, fill=LINE)
bullets(s, 0.72, 4.5, 11.6, 1.4, [
    [("Lattice-based Adaptor Signatures", {"bold": True, "color": TEXT, "size": 15}),
     ("  —  the first public implementation of an exotic post-quantum signature, ", {"color": MUTED, "size": 15}),
     ("demonstrated on a blockchain atomic swap", {"color": TEXT, "size": 15})],
    [("Reference: ", {"color": MUTED, "size": 15}),
     ("Esgin, Ersoy & Erkin, eprint 2020/845, Algorithm 2", {"color": ACCENT, "size": 15, "bold": True}),
     ("   ·   base: ", {"color": MUTED, "size": 15}),
     ("pq-crystals/dilithium", {"color": ACCENT, "size": 15, "bold": True})],
], size=15, gap=9, marker="—", marker_color=AMBER)
simple(s, 0.72, 6.5, 11.0, 0.5, "Royce Steven   ·   Supervisor: Wang Zhipeng   ·   2026",
       13, DIM, font=BODY)
notes(s, "This deck is a design defence: for every non-obvious implementation choice I made when "
         "building LAS on top of the Dilithium reference code, I state the question, what the paper's "
         "Algorithm 2 actually requires, the decision I took, and why it is correct and in-scope. "
         "The thesis novelty is twofold: first public LAS implementation, and first exotic PQ signature "
         "placed in a blockchain setting.")

# ----------------------------------------------------------------- 2 | FRAMING (2x2)
s = base(prs, "Where this work sits", "The empty cell: exotic × post-quantum", n=2)
# 2x2 quadrant on the right
qx, qy, qw, qh = 7.35, 1.95, 5.25, 4.45
cellw, cellh = qw/2, qh/2
labels = [  # (row, col, text, sub, color, filled)
    (0,0,"Basic · Classical","ECDSA, Schnorr", DIM, False),
    (0,1,"Exotic · Classical","ring, multisig,\nadaptor (Schnorr)", DIM, False),
    (1,0,"Basic · PQ","Dilithium, Falcon,\nSPHINCS+  (poqeth)", MUTED, False),
    (1,1,"Exotic · PQ","mostly paper-only", AMBER, True),
]
for (r,c,t,sub,col,fill) in labels:
    cx, cy = qx + c*cellw, qy + r*cellh
    card(s, cx+0.06, cy+0.06, cellw-0.12, cellh-0.12,
         fill=(CARD2 if fill else PANEL), line=(AMBER if fill else LINE))
    simple(s, cx+0.22, cy+0.24, cellw-0.4, 0.4, t, 13.5, (AMBER if fill else TEXT), font=BODY, bold=True)
    tb,tf = textbox(s, cx+0.22, cy+0.72, cellw-0.4, cellh-0.9)
    for i,row in enumerate(sub.split("\n")):
        p = para(tf, first=(i==0), space_after=1, line=1.0)
        _run(p, row, 11, col, font=BODY)
    if fill:
        simple(s, cx+0.22, cy+cellh-0.52, cellw-0.4, 0.4, "← LAS lands here",
               11.5, ACCENT, font=BODY, bold=True, italic=True)
simple(s, qx, qy-0.42, qw, 0.35, "The signature taxonomy", 12, MUTED, font=BODY, bold=True, spc=1.2)

bullets(s, 0.62, 2.05, 6.4, 4.4, [
    [("Blockchains sign with ECDSA / Schnorr; ", {"color":TEXT}),
     ("Shor's algorithm breaks both.", {"color":CORAL, "bold":True})],
    [("NIST standardised ", {"color":TEXT}),
     ("basic", {"color":TEXT,"bold":True}),
     (" PQ signatures — but ", {"color":TEXT}),
     ("exotic", {"color":AMBER,"bold":True}),
     (" ones (adaptor, ring, group) are mostly ", {"color":TEXT}),
     ("paper-only", {"color":AMBER,"bold":True}),
     (" in the PQ setting.", {"color":TEXT})],
    [("Adaptor signatures", {"color":ACCENT,"bold":True}),
     (" enable scriptless atomic swaps and payment channels.", {"color":TEXT})],
    [("Thesis = close that gap: a working LAS, ", {"color":TEXT}),
     ("benchmarked and demonstrated on-chain.", {"color":TEXT,"bold":True})],
], size=15, gap=12)
notes(s, "Context for why the project matters. The 2x2 framing is the spine of the whole thesis: the "
         "exotic-PQ cell is essentially empty in practice. Everything I defend later serves the goal of "
         "filling that cell with real, correct, benchmarked code rather than another paper.")

# ----------------------------------------------------------------- 3 | AGENDA
s = base(prs, "What I will defend", "Eight design decisions an examiner will question", n=3)
items = [
    ("1","Reuse Dilithium, don't reinvent","exotic = basic scheme + a few extra functions", ACCENT),
    ("2","Variant B (paper Alg. 2)","fold the statement into the Fiat–Shamir hash", ACCENT),
    ("3","The simplified scheme","drop Power2Round / hint vector / bit-decomposition", AMBER),
    ("4","The norm budget γ−κ−1","the one decision the whole adaptor hinges on", CORAL),
    ("5","Ternary witnesses","statement/witness pair = another key pair", ACCENT),
    ("6","q = 8380417 (≈2²³), not 2²⁴","reuse Dilithium's fixed NTT table", AMBER),
    ("7","Self-contained parameters","n=ℓ=4, κ=60, γ=κ·d·(n+ℓ)", ACCENT),
    ("8","Determinism + KATs","reproducibility & nonce-reuse safety", ACCENT),
]
cols, rows = 2, 4
cw, ch = 6.02, 1.12
x0, y0 = 0.62, 1.95
gx, gy = 0.18, 0.10
for idx,(num,t,sub,col) in enumerate(items):
    r, c = idx % rows, idx // rows
    cx, cy = x0 + c*(cw+gx), y0 + r*(ch+gy)
    card(s, cx, cy, cw, ch, fill=CARD, line=LINE)
    rect(s, cx, cy, 0.10, ch, fill=col)
    # number circle
    o = rect(s, cx+0.26, cy+ch/2-0.27, 0.54, 0.54, fill=PANEL, line=col, lw=1.5, shape=MSO_SHAPE.OVAL)
    tf=o.text_frame; _set_margins(tf,0.0); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    pp=tf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER; _run(pp,num,18,col,font=HEAD,bold=True)
    simple(s, cx+1.0, cy+0.16, cw-1.15, 0.4, t, 14.5, TEXT, font=BODY, bold=True)
    simple(s, cx+1.0, cy+0.6, cw-1.15, 0.4, sub, 11.5, MUTED, font=BODY, italic=True)
notes(s, "The roadmap. Each card is one slide. Decisions 1, 5, 7, 8 are 'reuse / faithful' choices in "
         "cyan; 3 and 6 are deliberate, supervisor-sanctioned deviations in amber; decision 4 — the "
         "norm budget — is in coral because it is the single point of failure of the whole construction.")

# ----------------------------------------------------------------- 4 | DECISION 1 reuse
s = base(prs, "Decision 1", "Reuse Dilithium's primitives — don't reinvent lattice arithmetic", n=4)
simple(s, 0.62, 1.78, 12.0, 0.5,
       "An exotic scheme = a basic scheme + a few extra algorithms.", 17, ACCENT, font=BODY, bold=True)
# left: what is reused vs added
card(s, 0.62, 2.4, 5.95, 3.9, fill=PANEL, line=LINE)
simple(s, 0.85, 2.55, 5.5, 0.35, "REUSED AS-IS (mode-independent)", 12, ACCENT, font=BODY, bold=True, spc=1.2)
bullets(s, 0.85, 2.98, 5.55, 1.7, [
    "NTT / inverse-NTT mod Q  (poly_ntt)",
    "SHAKE-128/256, Keccak  (fips202)",
    "Polynomial ring arithmetic mod (X^N + 1)",
    "Rejection / uniform sampling (poly_uniform)",
    "Challenge construction  (poly_challenge)",
], size=13, gap=5, marker="✓", marker_color=ACCENT)
simple(s, 0.85, 4.85, 5.5, 0.35, "ADDED ON TOP (new, self-contained)", 12, AMBER, font=BODY, bold=True, spc=1.2)
bullets(s, 0.85, 5.25, 5.55, 1.1, [
    "las_sign / las_verify  (FS-with-aborts core)",
    "PreSign · PreVerify · Adapt · Ext",
    "las_Amul, bound logic, serialisation",
], size=13, gap=5, marker="+", marker_color=AMBER)
# right: rationale + headline stat
bullets(s, 6.85, 2.4, 5.8, 2.7, [
    [("The NTT (table fixed to Q), Keccak and rejection sampling are the ", {"color":TEXT}),
     ("hard, audited, side-channel-sensitive", {"color":AMBER,"bold":True}),
     (" parts. Re-deriving them invites bugs and adds no thesis value.", {"color":TEXT})],
    [("LAS is layered as a ", {"color":TEXT}),
     ("self-contained scheme", {"color":ACCENT,"bold":True}),
     (" — its own dimensions and parameters — on those primitives.", {"color":TEXT})],
    [("Contribution must be a ", {"color":TEXT}),
     ("clean, visible diff", {"color":ACCENT,"bold":True}),
     (", not a fork that's impossible to audit.", {"color":TEXT})],
], size=14.5, gap=11)
card(s, 6.85, 5.25, 5.78, 1.05, fill=CARD2, line=ACCENT)
simple(s, 7.05, 5.36, 2.6, 0.8, "0", 50, ACCENT, font=HEAD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
tb,tf = textbox(s, 8.05, 5.36, 4.45, 0.85, anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True,space_after=0,line=1.05)
_run(p,"upstream Dilithium functions modified.\n",13.5,TEXT,font=BODY,bold=True)
p2=para(tf,space_after=0,line=1.0); _run(p2,"Verified in docs/FUNCTION_MAP.md — the report's reuse table.",11,MUTED,font=BODY)
notes(s, "The foundational decision. The key design fact of the thesis is that an exotic scheme is a basic "
         "scheme plus a handful of functions. So I reuse Dilithium's mode-independent internals and write "
         "only the adaptor layer. The headline is zero upstream functions modified — the entire "
         "contribution is additive, which makes it auditable and is exactly what the function-map "
         "deliverable documents.")

# ----------------------------------------------------------------- 5 | DECISION 2 variant B
s = base(prs, "Decision 2", "Variant B: fold the statement into the Fiat–Shamir hash", n=5)
# correction banner
card(s, 0.62, 1.78, 12.0, 0.62, fill=PANEL, line=AMBER)
tb,tf=textbox(s,0.8,1.78,11.7,0.62,anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True,space_after=0)
_run(p,"Correction after re-reading Algorithm 2:  ",12.5,AMBER,font=BODY,bold=True)
_run(p,"an earlier reading (“variant A”: z̃ = z+y, subtract Y at verify) was ",12.5,TEXT,font=BODY)
_run(p,"superseded",12.5,CORAL,font=BODY,bold=True)
_run(p,". The paper folds Y into the challenge.",12.5,TEXT,font=BODY)
# two code panels: Sign vs PreSign
code_panel(s, 0.62, 2.65, 5.95, 1.95, title="SIGN  (ordinary)", accent=ACCENT, lines=[
    "y <- S_gamma ;   w = A y",
    [("c = H(pk, ", {}), ("w", {"color":ACCENT,"bold":True}), (", M)", {})],
    "z = y + c*r ;   reject |z|inf > g-k",
])
code_panel(s, 6.7, 2.65, 5.95, 1.95, title="PRESIGN  (adaptor)", accent=AMBER, lines=[
    "y <- S_gamma ;   w = A y",
    [("c = H(pk, ", {}), ("w + Y", {"color":AMBER,"bold":True}), (", M)", {})],
    "z^ = y + c*r ;  reject |z^|inf > g-k-1",
])
simple(s, 0.62, 4.62, 12.0, 0.4,
       "The single algorithmic difference is the  “+ Y ”  inside the hash. That is the entire adaptor mechanism.",
       14, TEXT, font=BODY, italic=True)
# adapt / ext + payoff
card(s, 0.62, 5.15, 5.95, 1.2, fill=CARD, line=LINE)
bullets(s, 0.82, 5.26, 5.6, 1.05, [
    [("Adapt:", {"color":ACCENT,"bold":True}), ("  σ = (c, ẑ + y)   ", {"color":TEXT,"font":MONO}), ("— add the witness", {"color":MUTED})],
    [("Ext:", {"color":ACCENT,"bold":True}), ("    y = z − ẑ   ", {"color":TEXT,"font":MONO}), ("— subtract to recover it", {"color":MUTED})],
], size=13, gap=7, marker="")
card(s, 6.7, 5.15, 5.92, 1.2, fill=CARD2, line=ACCENT)
bullets(s, 6.9, 5.24, 5.6, 1.05, [
    [("The payoff: ", {"color":ACCENT,"bold":True}),
     ("the adapted σ is a ", {"color":TEXT}),
     ("fully ordinary signature", {"color":TEXT,"bold":True}),
     (".", {"color":TEXT})],
    [("Verify sees  ", {"color":TEXT}),
     ("Az − ct = w + Y", {"color":ACCENT,"font":MONO,"bold":True}),
     (",  matching c — the chain needs zero adaptor-awareness.", {"color":TEXT})],
], size=13, gap=6, marker="")
notes(s, "Why variant B and not the earlier reading. The defining move of the paper is that the statement Y "
         "is committed inside the Fiat-Shamir challenge: Sign hashes w, PreSign hashes w+Y. Adapt then adds "
         "the ternary witness and Ext subtracts. The crucial consequence is that an adapted signature is "
         "indistinguishable from an ordinary one — the on-chain verifier is plain Dilithium-style Verify, "
         "no special logic. I corrected an earlier 'variant A' reading after re-reading Algorithm 2 carefully.")

# ----------------------------------------------------------------- 6 | the 7 algorithms map (evidence)
s = base(prs, "Faithful to Algorithm 2", "Every paper step maps to one C function", n=6)
rows_data = [
    ("KeyGen",   "r ← S₁ ;  t = A·r",                 "las_keygen", ACCENT),
    ("Sign",     "c = H(pk, w, M) ;  z = y + c·r",       "las_sign", ACCENT),
    ("Verify",   "w' = A·z − c·t ;  c =? H(pk, w', M)", "las_verify", ACCENT),
    ("PreSign",  "c = H(pk, w+Y, M) ;  |ẑ|∞ ≤ γ−κ−1", "las_presign", AMBER),
    ("PreVerify","c =? H(pk, w'+Y, M)",                  "las_preverify", AMBER),
    ("Adapt",    "σ = (c, ẑ + y)",                       "las_adapt", AMBER),
    ("Ext",      "y = z − ẑ ;  check A·y = Y",        "las_ext", AMBER),
]
tx, ty0 = 0.62, 1.95
rowh = 0.62
# header
simple(s, tx+0.1, ty0, 2.3, 0.4, "ALGORITHM", 11, MUTED, font=BODY, bold=True, spc=1.2)
simple(s, tx+2.5, ty0, 5.6, 0.4, "PAPER STEP (Alg. 2)", 11, MUTED, font=BODY, bold=True, spc=1.2)
simple(s, tx+8.7, ty0, 3.4, 0.4, "C FUNCTION (ref/las.c)", 11, MUTED, font=BODY, bold=True, spc=1.2)
for i,(name,step,fn,col) in enumerate(rows_data):
    yy = ty0 + 0.42 + i*rowh
    card(s, tx, yy, 11.98, rowh-0.08, fill=(CARD if i%2==0 else PANEL), line=None)
    rect(s, tx, yy, 0.08, rowh-0.08, fill=col)
    simple(s, tx+0.28, yy, 2.2, rowh-0.08, name, 14, col, font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, tx+2.5, yy, 6.0, rowh-0.08, step, 13.5, TEXT, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, tx+8.7, yy, 3.2, rowh-0.08, fn+"()", 13, ACCENT, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "This is the evidence slide for fidelity. Each of the seven algorithms in the paper's Algorithm 2 "
         "maps to exactly one function in ref/las.c, and the full equation-by-equation correspondence is in "
         "docs/THEORY_IMPL_BRIDGE.md. Cyan rows are the base signature; amber rows are the four adaptor "
         "functions that constitute the exotic extension.")

# ----------------------------------------------------------------- 7 | DECISION 3 simplified scheme
s = base(prs, "Decision 3", "The simplified scheme: drop Power2Round, hint vector & bit-decomposition", n=7, accent=AMBER)
bullets(s, 0.62, 1.85, 6.15, 3.0, [
    [("Optimised Dilithium ", {"color":TEXT}),
     ("compresses t", {"color":AMBER,"bold":True}),
     (" (Power2Round) and ships a ", {"color":TEXT}),
     ("hint vector h", {"color":AMBER,"bold":True}),
     (" so the verifier can recover the high bits of w.", {"color":TEXT})],
    [("Those optimisations ", {"color":TEXT}),
     ("break the exact identity", {"color":CORAL,"bold":True}),
     (" the adaptor and witness-extraction depend on.", {"color":TEXT})],
    [("So I disable them — ", {"color":TEXT}),
     ("supervisor-sanctioned", {"color":ACCENT,"bold":True}),
     (": follow the paper's algorithm; a prototype may drop the reference's extra checks.", {"color":TEXT})],
], size=14.5, gap=12)
# the identity highlight
card(s, 0.62, 5.0, 6.15, 1.35, fill=CARD2, line=ACCENT)
simple(s, 0.82, 5.12, 5.8, 0.35, "THE IDENTITY THAT MUST HOLD EXACTLY", 11, ACCENT, font=BODY, bold=True, spc=1.0)
simple(s, 0.82, 5.5, 5.8, 0.55, "A·z − c·t  =  w + Y", 23, TEXT, font=MONO, bold=True)
simple(s, 0.82, 6.06, 5.8, 0.3, "holds only with no bit-dropping — Ext = z − ẑ must be exact",
       11.5, MUTED, font=BODY, italic=True)
# right: the trade-off (sizes)
simple(s, 7.0, 1.85, 5.6, 0.4, "THE PRICE WE PAY (and accept)", 12, AMBER, font=BODY, bold=True, spc=1.2)
sizes = [("LAS packed signature","4 672 B", CORAL),
         ("paper's optimised scheme","~3 210 B", MUTED),
         ("Dilithium-3 signature","3 309 B", MUTED)]
for i,(lab,val,col) in enumerate(sizes):
    yy = 2.35 + i*0.74
    card(s, 7.0, yy, 5.6, 0.64, fill=CARD, line=LINE)
    simple(s, 7.22, yy, 3.7, 0.64, lab, 13.5, TEXT, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, 10.7, yy, 1.75, 0.64, val, 16, col, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
bullets(s, 7.0, 4.75, 5.6, 1.6, [
    [("Cost is ", {"color":TEXT}), ("size only", {"color":AMBER,"bold":True}),
     (" — correctness and exact extraction are preserved.", {"color":TEXT})],
    [("Gains: ", {"color":ACCENT,"bold":True}),
     ("a clean algebraic spec, an auditable verifier, and a faithful match to the paper.", {"color":TEXT})],
], size=14, gap=10)
notes(s, "Decision 3 is the biggest deliberate deviation from optimised Dilithium. The hint vector exists "
         "purely to let the verifier reconstruct w's high bits after t is compressed. But the adaptor "
         "mechanism and the witness extractor rely on the exact identity Az-ct = w+Y; any bit-dropping "
         "destroys it. The supervisor explicitly said to follow the paper and simplify the reference's "
         "optimisations. The only cost is a larger signature — a size, not a correctness, penalty.")

# ----------------------------------------------------------------- 8 | DECISION 4 the norm budget (HERO)
s = base(prs, "Decision 4  ·  the critical one", "The norm budget: why PreSign rejects at γ−κ−1, not γ−κ", n=8, accent=CORAL)
simple(s, 0.62, 1.8, 12.0, 0.45,
       "This single off-by-one is what makes the whole adaptor correct. Loosen it and Verify rejects every adapted signature.",
       14.5, CORAL, font=BODY, bold=True)
# the inequality chain as a flow
chainy = 2.55
steps = [
    ("PreSign bound", "|ẑ|∞ ≤ γ−κ−1", AMBER),
    ("+ ternary witness", "|y|∞ ≤ 1", ACCENT),
    ("Adapt: z = ẑ + y", "|z|∞ ≤ γ−κ", ACCENT),
    ("✓ clears Verify", "bound γ−κ", ACCENT),
]
bw, bh, gap = 2.78, 1.35, 0.34
for i,(lab,val,col) in enumerate(steps):
    cx = 0.62 + i*(bw+gap)
    card(s, cx, chainy, bw, bh, fill=CARD, line=col)
    simple(s, cx+0.12, chainy+0.16, bw-0.24, 0.4, lab, 12, MUTED, font=BODY, bold=True)
    simple(s, cx+0.12, chainy+0.55, bw-0.24, 0.6, val, 19, col, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps)-1:
        simple(s, cx+bw+0.01, chainy, gap, bh, "→", 24, DIM, font=BODY, bold=True,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# the math
code_panel(s, 0.62, 4.25, 6.0, 1.65, title="THE BUDGET ARITHMETIC", accent=ACCENT, size=13, lines=[
    "|z|inf = |z^ + y|inf",
    "       <= |z^|inf + |y|inf",
    "       <= (g-k-1) + 1  =  g-k     (clears Verify)",
])
# the failure mode
card(s, 6.8, 4.25, 5.82, 1.65, fill=CARD, line=CORAL)
simple(s, 7.0, 4.38, 5.4, 0.35, "THE FAILURE MODE TO WATCH", 11.5, CORAL, font=BODY, bold=True, spc=1.0)
bullets(s, 7.0, 4.78, 5.45, 1.1, [
    [("Loosen PreSign to ", {"color":TEXT}), ("γ−κ", {"color":CORAL,"font":MONO,"bold":True}),
     (" and adapted z can reach ", {"color":TEXT}), ("γ−κ+1", {"color":CORAL,"font":MONO,"bold":True}),
     (" — over the bound.", {"color":TEXT})],
    [("Verify then rejects ", {"color":TEXT}), ("every", {"color":CORAL,"bold":True}),
     (" adapted signature. The budget, not packing, is the real constraint.", {"color":TEXT})],
], size=13, gap=7, marker="")
simple(s, 0.62, 6.1, 12.0, 0.4,
       "Generalises to AMHL:  PreSign bound γ−κ−K reserves budget K for a cumulative witness with |s_j|∞ ≤ j ≤ K.",
       12.5, MUTED, font=BODY, italic=True)
notes(s, "This is the decision the whole construction hinges on, and the failure mode the supervisor "
         "flagged. The witness is ternary, so its infinity norm is at most 1. Adapt forms z = z-hat + y. "
         "If PreSign already allowed z-hat up to gamma-kappa, then z could reach gamma-kappa+1, which "
         "exceeds the ordinary Verify bound and every adapted signature is rejected. By rejecting one "
         "tighter, at gamma-kappa-1, we leave exactly one unit of head-room for the witness. The AMHL "
         "multi-hop case generalises this to gamma-kappa-K.")

# ----------------------------------------------------------------- 9 | DECISION 5 ternary witnesses
s = base(prs, "Decision 5", "Ternary witnesses: the statement/witness pair is just another key pair", n=9)
code_panel(s, 0.62, 1.95, 5.95, 1.5, title="RELATION (hard under M-SIS / M-LWE)", accent=ACCENT, size=14, lines=[
    "y <- S_1^(n+l)      (ternary: -1, 0, 1)",
    "Y  =  A * y         (the statement)",
])
simple(s, 0.62, 3.65, 6.0, 0.45, "Identical in form to KeyGen:  r ← S₁,  t = A·r.",
       14.5, MUTED, font=BODY, italic=True)
# three reasons
reasons = [
    ("Norm head-room", "|y|∞ ≤ 1 is exactly the budget Adapt needs (Decision 4).", ACCENT),
    ("No separate Gen", "the hard relation R is literally the key-generation relation — one algorithm, reused.", ACCENT),
    ("Hiding", "M-SIS/M-LWE means knowing Y reveals nothing about y; publishing σ leaks y only to the pre-sig holder.", ACCENT),
]
for i,(t,d,col) in enumerate(reasons):
    yy = 4.4 + i*0.68
    card(s, 0.62, yy, 5.95, 0.6, fill=CARD, line=LINE)
    rect(s, 0.62, yy, 0.08, 0.6, fill=col)
    simple(s, 0.85, yy, 1.9, 0.6, t, 12.5, col, font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, 2.7, yy, 3.75, 0.6, d, 10.8, TEXT, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
# right: the on-chain leak intuition
card(s, 6.8, 1.95, 5.82, 3.78, fill=PANEL, line=LINE)
simple(s, 7.0, 2.1, 5.4, 0.4, "WHY SWAPS ARE ATOMIC", 12, AMBER, font=BODY, bold=True, spc=1.2)
bullets(s, 7.0, 2.55, 5.45, 3.7, [
    [("Alice pre-signs over statement Y; she alone knows witness y.", {"color":TEXT})],
    [("To claim, Alice ", {"color":TEXT}), ("adapts", {"color":ACCENT,"bold":True}),
     (" — publishing an ordinary σ = (c, ẑ+y) on chain A.", {"color":TEXT})],
    [("Bob holds the matching pre-sig σ̂ = (c, ẑ). He computes ", {"color":TEXT}),
     ("y = z − ẑ", {"color":ACCENT,"font":MONO,"bold":True}),
     (" — the witness is now public to him.", {"color":TEXT})],
    [("Bob uses y to complete the matching half on chain B. ", {"color":TEXT}),
     ("Both legs settle or neither does.", {"color":AMBER,"bold":True})],
], size=14, gap=14)
notes(s, "Decision 5 ties the cryptography to the application. Making the witness ternary serves three "
         "purposes at once: it provides exactly the norm head-room Adapt needs, it means the hard relation "
         "is identical to key generation so there is no separate setup, and M-SIS/M-LWE hardness hides y "
         "behind Y. The right panel is the atomicity argument: publishing the adapted signature on one "
         "chain mathematically leaks the witness to the counterparty, who completes the other leg.")

# ----------------------------------------------------------------- 10 | DECISION 6 modulus q
s = base(prs, "Decision 6", "q = 8 380 417 (≈2²³), not the paper's q ≈ 2²⁴", n=10, accent=AMBER)
bullets(s, 0.62, 1.9, 6.15, 2.7, [
    [("Reusing Dilithium's NTT ", {"color":TEXT}),
     ("pins the modulus", {"color":AMBER,"bold":True}),
     (": its root-of-unity table is fixed to Q = 8 380 417  (ROOT_OF_UNITY = 1753).", {"color":TEXT})],
    [("Changing q to 2²⁴ would require a ", {"color":TEXT}),
     ("new NTT table or schoolbook multiplication", {"color":AMBER,"bold":True}),
     (" — i.e. abandoning the core reuse.", {"color":TEXT})],
    [("Supervisor-sanctioned starting point: ", {"color":ACCENT,"bold":True}),
     ("start from the implementation's parameters, change nothing at first.", {"color":TEXT})],
], size=14.5, gap=12)
# correctness preserved box
card(s, 0.62, 4.75, 6.15, 1.6, fill=CARD2, line=ACCENT)
simple(s, 0.82, 4.88, 5.8, 0.35, "CORRECTNESS IS UNAFFECTED", 11.5, ACCENT, font=BODY, bold=True, spc=1.0)
simple(s, 0.82, 5.25, 5.8, 0.5, "Q  >  2γ      8 380 417  >  245 760", 18, TEXT, font=MONO, bold=True)
simple(s, 0.82, 5.78, 5.8, 0.5, "no wrap-around in the bound arithmetic ⇒ every identity still holds.",
       12, MUTED, font=BODY, italic=True)
# right: what changes / out of scope
card(s, 7.0, 1.9, 5.62, 4.45, fill=PANEL, line=LINE)
simple(s, 7.2, 2.05, 5.2, 0.4, "WHAT ACTUALLY CHANGES", 12, AMBER, font=BODY, bold=True, spc=1.2)
bullets(s, 7.2, 2.5, 5.25, 2.0, [
    [("Only the concrete ", {"color":TEXT}),
     ("M-SIS / M-LWE security margin", {"color":AMBER,"bold":True}),
     (" — not functionality.", {"color":TEXT})],
    [("Security analysis is ", {"color":TEXT}),
     ("explicitly out of scope", {"color":ACCENT,"bold":True}),
     (": treat hardness as a black box (supervisor).", {"color":TEXT})],
], size=14, gap=12)
card(s, 7.2, 4.7, 5.22, 1.45, fill=CARD, line=LINE)
simple(s, 7.4, 4.82, 4.9, 0.35, "DOCUMENTED FUTURE WORK", 11, MUTED, font=BODY, bold=True, spc=1.0)
simple(s, 7.4, 5.2, 4.9, 0.9,
       "Migration to q ≈ 2²⁴ is a separate, documented step with before/after benchmarks — optional, only if justified.",
       12.5, TEXT, font=BODY, anchor=MSO_ANCHOR.TOP, line=1.05)
notes(s, "Decision 6 is the second deliberate deviation, and it is forced by Decision 1. Dilithium's NTT "
         "hard-codes its root-of-unity table to Q = 8380417, so reusing the NTT pins the modulus. The "
         "paper wants ~2^24. Crucially this does not affect correctness because Q exceeds 2*gamma, so the "
         "norm arithmetic never wraps. The only thing that shifts is the concrete security margin, and "
         "security analysis is out of scope by the supervisor's instruction. Migration to 2^24 is logged "
         "as future work.")

# ----------------------------------------------------------------- 11 | DECISION 7 parameters + challenge
s = base(prs, "Decision 7", "Self-contained parameters & the challenge construction", n=11)
# parameter table
simple(s, 0.62, 1.85, 6.0, 0.4, "LAS PARAMETERS  (independent of DILITHIUM_MODE)", 12, ACCENT, font=BODY, bold=True, spc=1.0)
param_rows = [
    ("N = d", "256", "ring degree, R_q = Z_q[X]/(X^N+1)"),
    ("n = l", "4", "module dimension / extra columns"),
    ("n+l", "8", "length of r, y, z"),
    ("kappa", "60", "challenge weight  ||c||_1"),
    ("gamma", "122 880", "= kappa*d*(n+l) = 60*256*8"),
]
ty0 = 2.32
for i,(sym,val,desc) in enumerate(param_rows):
    yy = ty0 + i*0.66
    card(s, 0.62, yy, 6.0, 0.58, fill=(CARD if i%2==0 else PANEL))
    simple(s, 0.8, yy, 1.5, 0.58, sym, 14, AMBER, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, 2.25, yy, 1.5, 0.58, val, 14, TEXT, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, 3.8, yy, 2.75, 0.58, desc, 10.3, MUTED, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
# right: challenge + acceptance
code_panel(s, 6.85, 1.95, 5.78, 1.55, title="CHALLENGE  c = H(...)  (Dilithium poly_challenge)", accent=ACCENT, size=13.5, lines=[
    "||c||_1 = kappa = 60   (exactly 60 nonzero)",
    "||c||_inf = 1          (entries are +-1)",
])
card(s, 6.85, 3.7, 5.78, 2.65, fill=PANEL, line=LINE)
simple(s, 7.05, 3.83, 5.4, 0.4, "ACCEPTANCE RATE — MEASURED DIRECTLY", 12, ACCENT, font=BODY, bold=True, spc=1.0)
simple(s, 7.05, 4.3, 5.4, 0.85, "≈ 37%", 44, ACCENT, font=HEAD, bold=True)
simple(s, 9.1, 4.45, 3.4, 0.6, "per attempt\n≈ 2.7 attempts / sig", 13, TEXT, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
simple(s, 7.05, 5.35, 5.4, 0.9,
       "matches  (1 − κ/γ)^((n+ℓ)·N)  ≈  e^(−1)  ≈  36.8%,\nvia the las_attempts counter (not a timing-ratio estimate).",
       12, MUTED, font=BODY, italic=True, line=1.1)
notes(s, "Decision 7 documents the concrete parameters and that they are self-contained — the scheme is "
         "mode-independent, built under MODE=3 but not depending on it. Gamma is chosen as kappa*d*(n+l) to "
         "make the M-SIS hardness parameter adequate, not to maximise acceptance. The challenge reuses "
         "Dilithium's construction with weight kappa=60. I measure the rejection-sampling acceptance "
         "directly with an attempt counter — about 37%, matching the closed form e^-1 — rather than "
         "estimating it from a timing ratio, which earlier gave a biased 23%.")

# ----------------------------------------------------------------- 12 | DECISION 8 determinism + KAT
s = base(prs, "Decision 8", "Determinism & known-answer tests — reproducibility by construction", n=12)
code_panel(s, 0.62, 1.95, 6.05, 1.7, title="DERIVED MASK SEED (standard FS derandomisation)", accent=ACCENT, size=13, lines=[
    "seed = SHAKE256( tag || sk || [Y] || M )",
    "  tag 0x00 = Sign      tag 0x01 = PreSign (binds Y)",
    "shared sign_core / presign_core with random path",
])
bullets(s, 0.62, 3.95, 6.05, 2.4, [
    [("Same distribution and validity as the randomised path — ", {"color":TEXT}),
     ("not a change to the scheme", {"color":ACCENT,"bold":True}),
     (", just derandomisation (as in deterministic Dilithium).", {"color":TEXT})],
    [("Removes the ", {"color":TEXT}),
     ("nonce-reuse failure mode", {"color":AMBER,"bold":True}),
     (" (no fresh per-signature RNG to mishandle).", {"color":TEXT})],
    [("Enables ", {"color":TEXT}),
     ("byte-identical KATs", {"color":ACCENT,"bold":True}),
     (" across runs and machines — a cross-check anchor for any on-chain verifier.", {"color":TEXT})],
], size=14, gap=13)
# right: robustness evidence cards
ev = [
    ("1000 × 3", "iterations, modes 2/3/5 — 100% correct", ACCENT),
    ("4672 / 4672", "single-byte signature flips — all rejected", ACCENT),
    ("pinned digest", "SHAKE256 KAT matches:  f7fc40f0b775…", ACCENT),
]
for i,(big,desc,col) in enumerate(ev):
    yy = 1.95 + i*1.5
    card(s, 6.95, yy, 5.68, 1.32, fill=CARD, line=col)
    simple(s, 7.2, yy+0.16, 5.2, 0.6, big, 26, col, font=HEAD, bold=True)
    simple(s, 7.2, yy+0.8, 5.25, 0.45, desc, 13, TEXT, font=BODY)
notes(s, "Decision 8 is engineering rigour. The per-signature mask is derived as a SHAKE256 hash of a "
         "domain tag, the secret key, optionally the statement, and the message. This is standard "
         "Fiat-Shamir derandomisation — identical distribution, no scheme change — but it removes the "
         "nonce-reuse risk and makes the whole pipeline reproducible, which is what lets me pin a "
         "known-answer-test digest. The right column is the robustness evidence: 1000 iterations across "
         "three modes, every one of the 4672 byte-flips rejected, and a pinned KAT digest.")

# ----------------------------------------------------------------- 13 | how I know it's right (verification)
s = base(prs, "How I know it's faithful", "Test assertions encode the paper's correctness theorems", n=13)
test_rows = [
    ("PreVerify accepts the pre-sig", "pre-signature correctness", ACCENT),
    ("Verify REJECTS the pre-sig  (tripwire)", "statement binding — +Y really is in the hash", AMBER),
    ("Verify accepts the adapted σ", "pre-signature adaptability", ACCENT),
    ("Ext recovers y, and A·y = Y exactly", "witness extractability", ACCENT),
    ("Adapt hop 1 with s_K ⇒ Verify fails", "AMHL wormhole resistance", AMBER),
    ("Flip any bit ⇒ Verify fails", "basic unforgeability", ACCENT),
]
ty0=1.95; rh=0.66
for i,(claim,thm,col) in enumerate(test_rows):
    yy=ty0+i*rh
    card(s, 0.62, yy, 7.5, rh-0.08, fill=(CARD if i%2==0 else PANEL))
    rect(s, 0.62, yy, 0.08, rh-0.08, fill=col)
    o=rect(s, 0.82, yy+(rh-0.08)/2-0.13, 0.26, 0.26, fill=PANEL, line=col, lw=1.5, shape=MSO_SHAPE.OVAL)
    tf=o.text_frame;_set_margins(tf,0.0);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    pp=tf.paragraphs[0];pp.alignment=PP_ALIGN.CENTER;_run(pp,"✓",10,col,font=BODY,bold=True)
    simple(s, 1.2, yy, 4.45, rh-0.08, claim, 12.5, TEXT, font=BODY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    simple(s, 5.65, yy, 2.4, rh-0.08, thm, 10.3, MUTED, font=BODY, italic=True, anchor=MSO_ANCHOR.MIDDLE)
# right: the headline benchmark result
card(s, 8.35, 1.95, 4.27, 3.85, fill=CARD2, line=ACCENT)
simple(s, 8.55, 2.1, 3.9, 0.4, "THE COUNTER-INTUITIVE RESULT", 11.5, ACCENT, font=BODY, bold=True, spc=1.0)
simple(s, 8.55, 2.55, 3.9, 0.7, "≈ 0", 42, ACCENT, font=HEAD, bold=True)
simple(s, 8.55, 3.35, 3.9, 0.7, "adaptor overhead:\nPreSign ≈ Sign,  PreVerify ≈ Verify", 13, TEXT, font=BODY, line=1.1)
rect(s, 8.55, 4.35, 3.85, 0.02, fill=LINE)
simple(s, 8.55, 4.5, 3.9, 1.2, "vs the classical adaptor's ~4× DLEQ overhead. The PQ price is communication (×29–89 sizes), not computation.",
       12.5, MUTED, font=BODY, italic=True, line=1.12)
notes(s, "How I demonstrate fidelity rather than just claim it. Each property the paper proves is encoded "
         "as a hard assertion that runs every iteration. The most important is the 'tripwire': feeding a "
         "pre-signature to ordinary Verify must FAIL, which proves the +Y term is genuinely inside the "
         "hash and the pre-sig is unspendable. The right panel is the headline evaluation finding: LAS's "
         "adaptor overhead is essentially zero, unlike the classical adaptor's ~4x DLEQ cost — the "
         "post-quantum price is paid in size, not speed.")

# ----------------------------------------------------------------- 14 | deviations register (honesty)
s = base(prs, "Honest accounting", "Every deviation from the paper, and its impact", n=14, accent=AMBER)
hdr = ["PROPERTY", "PAPER", "THIS WORK", "IMPACT"]
data = [
    ("Modulus q", "≈ 2²⁴", "8 380 417 (≈2²³)", "security margin only — Q>2γ keeps correctness"),
    ("Hint vector h", "used (optimised)", "disabled (simplified)", "larger sig (4672 B); exact identity preserved"),
    ("Signature packing", "~3 210 B", "4 672 B (measured)", "size only; validating decoder added"),
    ("Extracted witness", "may carry noise", "exact (y = z − ẑ)", "the “knowledge gap” — noted as future work"),
    ("Multi-hop PCN", "AMHL, γ−κ−K", "AMHL implemented", "functionally matches; privacy variant is future"),
]
tx, ty0 = 0.62, 1.95
colx = [0.0, 2.2, 4.25, 6.6]
colw = [2.2, 2.05, 2.35, 5.4]
# header row
rect(s, tx, ty0, 11.98, 0.5, fill=CARD2)
for j,h in enumerate(hdr):
    simple(s, tx+colx[j]+0.15, ty0, colw[j]-0.2, 0.5, h, 11.5,
           (AMBER if j>=2 else MUTED), font=BODY, bold=True, spc=1.0, anchor=MSO_ANCHOR.MIDDLE)
rh=0.78
for i,row in enumerate(data):
    yy=ty0+0.5+i*rh
    rect(s, tx, yy, 11.98, rh, fill=(CARD if i%2==0 else PANEL))
    for j,val in enumerate(row):
        col = TEXT if j<2 else (ACCENT if j==2 else MUTED)
        bold = (j==0)
        fnt = MONO if (j==1 or (j==2 and i<3)) else BODY
        simple(s, tx+colx[j]+0.15, yy, colw[j]-0.25, rh, val, 11.5 if j<3 else 11,
               (TEXT if j==0 else col), font=fnt, bold=bold, anchor=MSO_ANCHOR.MIDDLE, line=1.0)
simple(s, 0.62, 6.55, 12.0, 0.4,
       "Stated plainly because examiners reward honest critique — not because any of these breaks the construction.",
       12.5, MUTED, font=BODY, italic=True)
notes(s, "The honesty slide. Distinction-level work shows it knows its own limits. Four deliberate "
         "deviations — the modulus, the dropped hint vector, the larger packed size, and the exact "
         "(rather than noisy) witness — each with its real impact. The 'knowledge gap', that our extracted "
         "witness is exact whereas the paper's relaxed setting allows noise that grows over long payment "
         "chains, is the one I flag explicitly as future work. None of these breaks correctness.")

# ----------------------------------------------------------------- 15 | CONCLUSION
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
rect(s, 0.0, 0.0, 0.18, SH, fill=ACCENT)
dot_grid(s, 11.55, 0.45, 4, 4, step=0.34, r=0.05, color=LATTICE)
simple(s, 0.7, 0.7, 11.0, 0.35, "DEFENCE SUMMARY", 13, ACCENT, font=BODY, bold=True, spc=2.6)
simple(s, 0.68, 1.12, 11.8, 0.9, "Every choice traces to three principles", 34, TEXT, font=HEAD, bold=True)
princ = [
    ("Faithful to Algorithm 2", "variant B, the +Y-in-hash mechanism, the γ−κ−1 budget, the seven functions — mapped equation-by-equation.", ACCENT),
    ("Maximal reuse of audited primitives", "NTT, SHAKE, sampling reused as-is; 0 upstream functions modified; the contribution is a clean, additive diff.", ACCENT),
    ("Every deviation justified & in-scope", "simplified scheme and q=2²³ are deliberate, supervisor-sanctioned, correctness-preserving; security analysis is out of scope.", AMBER),
]
for i,(t,d,col) in enumerate(princ):
    yy = 2.35 + i*1.18
    card(s, 0.7, yy, 11.9, 1.05, fill=CARD, line=None)
    rect(s, 0.7, yy, 0.10, 1.05, fill=col)
    o=rect(s, 0.95, yy+1.05/2-0.30, 0.6, 0.6, fill=PANEL, line=col, lw=1.5, shape=MSO_SHAPE.OVAL)
    tf=o.text_frame;_set_margins(tf,0.0);tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    pp=tf.paragraphs[0];pp.alignment=PP_ALIGN.CENTER;_run(pp,str(i+1),20,col,font=HEAD,bold=True)
    simple(s, 1.75, yy+0.14, 10.6, 0.4, t, 17, TEXT, font=BODY, bold=True)
    simple(s, 1.75, yy+0.55, 10.6, 0.45, d, 12.5, MUTED, font=BODY, line=1.05)
simple(s, 0.7, 6.15, 11.9, 0.7,
       "Result: the first public LAS implementation — correct, reproducible, benchmarked against two baselines, and demonstrated on-chain.",
       14, ACCENT, font=BODY, bold=True, italic=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 15)
notes(s, "The closing argument. Whatever an examiner asks, the answer reduces to one of three principles: "
         "the implementation is faithful to the paper's Algorithm 2; it reuses the audited Dilithium "
         "primitives maximally with zero upstream modifications; and every deviation is deliberate, "
         "justified, and within the agreed scope. The outcome is the first working, reproducible, "
         "benchmarked LAS, demonstrated in a blockchain atomic swap — filling the empty exotic-PQ cell.")

# ----------------------------------------------------------------- write
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LAS_Design_Defence.pptx")
prs.save(out)
print("SAVED:", out)
print("slides:", len(prs.slides._sldIdLst))
